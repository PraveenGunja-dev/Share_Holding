"""
=============================================================================
  Weekly Shareholder Movement â†’ PowerPoint Report Generator
  -----------------------------------------------------------
  Reads shareholder data from PostgreSQL tables and generates a professional
  PowerPoint presentation matching the Adani template structure.

  Slide Structure (after removing slides 3, 7, 12, 13):
    1.  Title Slide (from template)
    2.  Table of Contents (rebuilt)
    3.  Top 20 Institutional Shareholders
    4.  Top 20 Buyers
    5.  Top 20 Sellers
    6.  Top 10 FII's & FPI's
    7.  Top 10 MF's  (Active + Passive sub-tables)
    8.  Top 10 Insurance & PF's
    9.  Top 10 AIF's
   10.  Thank You (from template)

  Usage:
      1. Update config.py with your DB credentials and table names.
      2. pip install -r requirements.txt
      3. python generate_ppt.py
=============================================================================
"""

import sys
import os
from pathlib import Path
import copy
import math
from datetime import datetime
from decimal import Decimal, InvalidOperation, ROUND_DOWN

import sqlite3
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
import re
from typing import Optional

# Load environment variables from the project .env file so scripts run directly
_here = Path(__file__).resolve().parent
_env_candidates = [
    _here / ".env",                # root-level .env (current file directory)
    _here.parent / ".env",         # parent folder .env (legacy layout)
]
try:
    from dotenv import load_dotenv  # type: ignore
    for _env_path in _env_candidates:
        if _env_path.exists():
            load_dotenv(dotenv_path=_env_path)
            break
except Exception:
    # Fallback: minimal .env reader for KEY=VALUE or KEY="VALUE"
    try:
        for _env_path in _env_candidates:
            if _env_path.exists():
                for line in _env_path.read_text(encoding="utf-8").splitlines():
                    s = line.strip()
                    if not s or s.startswith("#") or "=" not in s:
                        continue
                    k, v = s.split("=", 1)
                    k = k.strip()
                    v = v.strip().strip("'\"")
                    # Do not override already-set env vars
                    if k and k not in os.environ:
                        os.environ[k] = v
                break
    except Exception:
        pass

def _get_env(key: str, default: str = "") -> str:
    v = os.getenv(key)
    return v if v is not None and str(v).strip() != "" else default


def _parse_date_col(name: str) -> Optional[datetime]:
    s = str(name or "").strip()
    for fmt in (
        "%d/%b/%y",
        "%d-%b-%y",
        "%d/%b/%Y",
        "%m/%d/%Y",
        "%m/%d/%y",
        "%d/%m/%Y",
        "%d/%m/%y",
    ):
        try:
            return datetime.strptime(s, fmt)
        except Exception:
            pass
    return None


def _get_bu_id() -> int:
    raw = (_get_env("WSHP_BU_ID", "1") or "1").strip()
    try:
        return int(raw)
    except Exception:
        return 1


def _get_date_range_label() -> Optional[str]:
    env_date1 = (_get_env("WSHP_DATE1", "") or "").strip()
    env_date2 = (_get_env("WSHP_DATE2", "") or "").strip()
    if env_date1 and env_date2:
        d1 = _parse_date_col(env_date1)
        d2 = _parse_date_col(env_date2)
        if d1 and d2:
            # Match analysis convention: previous vs current
            return f"{d2.strftime('%d-%b-%y')} vs {d1.strftime('%d-%b-%y')}"

    if PREVIOUS_WEEK_DATE and CURRENT_WEEK_DATE:
        return f"{PREVIOUS_WEEK_DATE} vs {CURRENT_WEEK_DATE}"
    return None

# Load DB config from environment
_default_db_path = str((_here.parent.parent / "WeeklyShareHolding_Update7.db").resolve())
DB_PATH = _get_env("DB_PATH", _default_db_path).strip()
if not DB_PATH:
    raise ValueError("Missing required DB env var: DB_PATH")

# Schema
SCHEMA_NAME = _get_env("WSHP_SCHEMA", "public")

# Table names mapped from existing env vars created by analysis generator
TABLE_NAMES = {
    # Institutional Top 20 (overall)
    "top_20_institutional": _get_env("WSHP_TOP20_TABLE", "Top 20 Holders"),
    # Buyers/Sellers sections (updated to Top 20)
    "top_20_buyers": _get_env("WSHP_TOP20_BUYERS_TABLE", "Top 20 Buyers"),
    "top_20_sellers": _get_env("WSHP_TOP20_SELLERS_TABLE", "Top 20 Sellers"),
    # FII/FPI block (use Top 20 FII table)
    "top_10_fii_fpi": _get_env("WSHP_TOP20_FII_TABLE", "Top 20 Holders FII"),
    # MF Active/Passive
    "top_10_mf_active": _get_env("WSHP_TOP20_MF_TABLE", "Top 20 Active Holders MF"),
    "top_10_mf_passive": _get_env("WSHP_TOP20_PASSIVE_MF_TABLE", "Top 20 Holder Passive MF"),
    # Insurance/PF and AIF
    "top_10_insurance_pf": _get_env("WSHP_TOP20_INS_PF_TABLE", "Top 20 Holders INS PF"),
    "top_10_aif": _get_env("WSHP_TOP20_AIF_TABLE", "Top 20 Holders AIF"),

    # New Entry / Exits
    "entry": _get_env("WSHP_ENTRY_TABLE", "Entry"),
    "exit": _get_env("WSHP_EXIT_TABLE", "Exit"),
}

# Report settings
PORTFOLIO_NAME = _get_env("WSHP_COMPANY_NAME", "Adani Portfolio")
CURRENT_WEEK_DATE = _get_env("WSHP_CURRENT_WEEK_DATE", "")
PREVIOUS_WEEK_DATE = _get_env("WSHP_PREVIOUS_WEEK_DATE", "")

# Attempt to derive dates dynamically from the analysis table when not supplied
ANALYSIS_TABLE = _get_env("WSHP_ANALYSIS_TABLE", "analysis")
def _derive_report_dates_from_db():
    try:
        conn = sqlite3.connect(DB_PATH)
        cur = conn.cursor()
        cur.execute(f'PRAGMA table_info("{ANALYSIS_TABLE}")')
        cols = [r[1] for r in cur.fetchall()]
        conn.close()
        import re
        date_cols = []
        for c in cols:
            if re.match(r"^\d{1,2}/\d{1,2}/\d{2,4}$", c.strip()):
                from datetime import datetime as _dt
                for fmt in ("%m/%d/%Y", "%m/%d/%y", "%d/%m/%Y", "%d/%m/%y"):
                    try:
                        dt = _dt.strptime(c.strip(), fmt)
                        date_cols.append((dt, c.strip()))
                        break
                    except Exception:
                        continue
        if len(date_cols) < 2:
            return None, None
        date_cols.sort(key=lambda x: x[0], reverse=True)
        latest, prev = date_cols[0][0], date_cols[1][0]
        return latest.strftime("%d-%b-%y"), prev.strftime("%d-%b-%y")
    except Exception:
        return None, None

if not CURRENT_WEEK_DATE or not PREVIOUS_WEEK_DATE:
    _d1, _d2 = _derive_report_dates_from_db()
    if _d1 and _d2:
        CURRENT_WEEK_DATE = CURRENT_WEEK_DATE or _d1
        PREVIOUS_WEEK_DATE = PREVIOUS_WEEK_DATE or _d2

# PowerPoint I/O
TEMPLATE_PPT = _get_env("WSHP_PPT_TEMPLATE", "Weekly Shareholder Movement_Template.pptx")
OUTPUT_PPT_FILENAME = _get_env("WSHP_PPT_OUTPUT", "Weekly_ShareHolding_Report.pptx")

# Slide dimensions
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5


# â”€â”€ Color Palette (matching template) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)   # White
DATA_FONT_COLOR   = RGBColor(0x00, 0x00, 0x00)   # Black
GREEN_COLOR       = RGBColor(0x00, 0xB0, 0x50)   # Green for positive changes
RED_COLOR         = RGBColor(0xFF, 0x00, 0x00)    # Red for negative changes
WHITE_BG          = RGBColor(0xFF, 0xFF, 0xFF)    # White background
HEADER_BG_COLOR   = RGBColor(0x00, 0x2B, 0x5C)    # Dark navy header background

# Font sizes (in EMUs from template: 127000 = Pt(10))
HEADER_FONT_SIZE = Pt(10)
DATA_FONT_SIZE   = Pt(10)
TITLE_FONT_SIZE  = Pt(16)   # ~203200 EMU
NOTE_FONT_SIZE   = Pt(9)    # ~114300 EMU
FOOTER_FONT_SIZE = Pt(8)    # ~101600 EMU


# â”€â”€ Database Functions â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def get_db_connection():
    """Establish and return a SQLite database connection."""
    try:
        conn = sqlite3.connect(DB_PATH)
        return conn
    except Exception as e:
        print(f"Database connection failed: {e}")
        sys.exit(1)


def fetch_table_data(conn, table_name: str, schema: str = "public") -> pd.DataFrame:
    """Fetch all rows from a table and return as a Pandas DataFrame."""
    query = f'SELECT * FROM "{table_name}"'
    try:
        df = pd.read_sql_query(query, conn)
        return df
    except Exception as e:
        print(f"   âš ï¸  Error reading table '{table_name}': {e}")
        return pd.DataFrame()


def _pick_latest_daterange(conn, table_name: str, schema: str = "public", bu_id: Optional[int] = None) -> Optional[str]:
    """Return the latest DateRange value for a table if the column exists."""
    try:
        cur = conn.cursor()
        cur.execute(f'PRAGMA table_info("{table_name}")')
        cols = [r[1] for r in cur.fetchall()]
        if "DateRange" not in cols:
            return None

        resolved_bu_id = bu_id if bu_id is not None else _get_bu_id()
        where_bu = ""
        if "bu_id" in cols:
            where_bu = f"WHERE bu_id = {int(resolved_bu_id)}"

        cur.execute(
            f"""
            SELECT DateRange
            FROM "{table_name}"
            {where_bu}
            ORDER BY rowid DESC
            LIMIT 1
            """
        )
        row = cur.fetchone()
        return row[0] if row else None
    except Exception:
        return None


def fetch_table_data_latest(
    conn,
    table_name: str,
    schema: str = "public",
    date_range: Optional[str] = None,
    bu_id: Optional[int] = None,
) -> pd.DataFrame:
    """Fetch a top-table snapshot.

    If DateRange column exists:
    - use the provided date_range, else auto-pick latest
    - filter rows to that DateRange
    Otherwise, returns full table.
    """
    df = fetch_table_data(conn, table_name, schema)
    if df is None or df.empty:
        return df

    if bu_id is not None and any(str(c) == "bu_id" for c in df.columns):
        try:
            df = df[df["bu_id"].astype("Int64") == int(bu_id)].copy()
        except Exception:
            try:
                df = df[df["bu_id"].astype(str) == str(int(bu_id))].copy()
            except Exception:
                pass

    if not any(str(c) == "DateRange" for c in df.columns):
        return df

    picked = date_range or _pick_latest_daterange(conn, table_name, schema, bu_id=bu_id)
    if not picked:
        return df
    try:
        return df[df["DateRange"].astype(str) == str(picked)].copy()
    except Exception:
        return df


# â”€â”€ Helper: Set cell styling â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def set_cell_style(cell, text, font_size=DATA_FONT_SIZE, bold=False,
                   font_color=DATA_FONT_COLOR, bg_color=None,
                   alignment=PP_ALIGN.CENTER, font_name="Arial"):
    """Apply consistent styling to a table cell."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    para = tf.paragraphs[0]
    para.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    run = para.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color
    run.font.name = font_name

    # Compact margins
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)

    if bg_color:
        try:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        except Exception:
            pass


def set_cell_border(cell, border_color="D9D9D9", width="6350"):
    """Add thin borders to a cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ["lnL", "lnR", "lnT", "lnB"]:
        ln = tcPr.makeelement(qn(f"a:{edge}"), {"w": width, "cap": "flat", "cmpd": "sng"})
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": border_color})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        # Add prstDash
        prstDash = ln.makeelement(qn("a:prstDash"), {"val": "solid"})
        ln.append(prstDash)
        tcPr.append(ln)


def remove_table_styling(table):
    """Remove default theme styling from the table."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = tbl.makeelement(qn("a:tblPr"), {})
        tbl.insert(0, tblPr)
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")


def merge_cells(table, start_row, start_col, end_row, end_col):
    """Merge cells in a table from (start_row, start_col) to (end_row, end_col)."""
    table.cell(start_row, start_col).merge(table.cell(end_row, end_col))


def get_change_color(value_str):
    """Return color based on positive/negative change value."""
    try:
        val_str = str(value_str).strip().replace(",", "").replace("(", "-").replace(")", "")
        val = float(val_str)
        if val > 0:
            return GREEN_COLOR
        elif val < 0:
            return RED_COLOR
        else:
            return DATA_FONT_COLOR
    except (ValueError, TypeError):
        if "-" in str(value_str).strip() and str(value_str).strip() != "-":
            return RED_COLOR
        return DATA_FONT_COLOR


def _format_int(value):
    """Best-effort conversion of a numeric value to an integer string (no decimals).
    - Handles strings with commas and parentheses negatives like (1,234.56)
    - Returns empty string for None/blank inputs
    - Leaves non-numeric strings unchanged
    """
    try:
        if value is None:
            return ""
        s = str(value).strip()
        if s == "":
            return ""
        # Normalize
        neg = False
        if s.startswith("(") and s.endswith(")"):
            neg = True
            s = s[1:-1]
        s = s.replace(",", "")
        num = float(s)
        ival = int(round(num))
        if neg:
            ival = -ival
        return str(ival)
    except Exception:
        return str(value if value is not None else "")


def _format_holding(value):
    """Best-effort formatting for holding values with 2 decimals.
    - Handles strings with commas and parentheses negatives like (1,234.56)
    - Returns empty string for None/blank inputs
    - Leaves non-numeric strings unchanged
    """
    try:
        if value is None:
            return ""
        s = str(value).strip()
        if s == "":
            return ""
        neg = False
        if s.startswith("(") and s.endswith(")"):
            neg = True
            s = s[1:-1].strip()
        s = s.replace(",", "")
        num = float(s)
        if neg:
            num = -num
        return f"{num:.2f}"
    except Exception:
        return str(value if value is not None else "")


def _format_change_value(value):
    """Format change values to 2 decimals using decimal-safe truncation toward zero."""
    try:
        if value is None:
            return ""
        s = str(value).strip()
        if s == "":
            return ""
        neg = False
        if s.startswith("(") and s.endswith(")"):
            neg = True
            s = s[1:-1].strip()
        s = s.replace(",", "")
        num = Decimal(s)
        if neg:
            num = -num
        q = num.quantize(Decimal("0.00"), rounding=ROUND_DOWN)
        return f"{q:.2f}"
    except (InvalidOperation, ValueError, TypeError):
        return str(value if value is not None else "")


def _format_pct(value):
    """Best-effort formatting for percentage values with a trailing '%' sign.
    - Returns empty string for None/blank inputs
    - Preserves existing '%' suffix if already present
    - Handles strings with commas and parentheses negatives
    - Leaves non-numeric strings unchanged except for appending '%'
    """
    try:
        if value is None:
            return ""
        s = str(value).strip()
        if s == "":
            return ""
        if s.endswith("%"):
            return s
        neg = False
        if s.startswith("(") and s.endswith(")"):
            neg = True
            s = s[1:-1].strip()
        s = s.replace(",", "")
        num = float(s)
        out = f"{num:g}%"
        if neg and not out.startswith("-"):
            out = f"-{out}"
        return out
    except Exception:
        s = "" if value is None else str(value).strip()
        if not s:
            return ""
        return s if s.endswith("%") else f"{s}%"


# â”€â”€ Slide Builders â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def _ensure_table_data_row_capacity(table, data_rows_needed: int, header_rows: int = 2):
    """Ensure the table has exactly header_rows + data_rows_needed rows by
    appending or removing rows at the bottom, preserving styling by cloning
    the last row when adding.
    """
    try:
        total_needed = max(header_rows, int(header_rows + max(0, int(data_rows_needed))))
    except Exception:
        total_needed = header_rows
    try:
        # Expand by cloning last row
        while len(table.rows) < total_needed:
            try:
                tbl = table._tbl
                last_tr = tbl.tr_lst[-1]
                new_tr = copy.deepcopy(last_tr)
                tbl.append(new_tr)
            except Exception:
                break
        # Shrink by removing rows from bottom (but never remove header rows)
        while len(table.rows) > total_needed:
            try:
                tbl = table._tbl
                tr_list = tbl.tr_lst
                if len(tr_list) <= header_rows:
                    break
                tbl.remove(tr_list[-1])
            except Exception:
                break
    except Exception:
        pass

def add_blank_slide(prs):
    """Add a blank slide using layout index 6."""
    slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(slide_layout)

def _delete_slides(prs, idxs):
    """Delete slides by indices (descending) from a Presentation."""
    try:
        sldIdLst = prs.slides._sldIdLst
        slide_ids = list(sldIdLst)
        for i in sorted(set([i for i in idxs if 0 <= i < len(slide_ids)]), reverse=True):
            sldIdLst.remove(slide_ids[i])
    except Exception:
        pass


def _find_slide_by_title(prs, title_text: str):
    """Find a slide whose TITLE placeholder text matches title_text (case-insensitive)."""
    if not title_text:
        return None
    def _norm(s: str) -> str:
        return " ".join(str(s or "").split()).strip().lower()
    want = _norm(title_text)
    for slide in prs.slides:
        try:
            for shape in slide.shapes:
                if getattr(shape, "is_placeholder", False):
                    phf = shape.placeholder_format
                    if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shape, "text_frame"):
                        txt = shape.text_frame.text or ""
                        txtn = _norm(txt)
                        if txtn == want or want in txtn:
                            return slide
        except Exception:
            continue
    return None

def _table_header_signature(tbl) -> str:
    """Return a normalized signature string for the first two rows and first 8 columns of a table."""
    parts = []
    try:
        rows = min(2, len(tbl.rows))
        cols = min(8, len(tbl.columns))
        for r in range(rows):
            row_parts = []
            for c in range(cols):
                try:
                    t = (tbl.cell(r, c).text or "").strip().lower()
                except Exception:
                    t = ""
                row_parts.append(" ".join(t.split()))
            parts.append("|".join(row_parts))
    except Exception:
        return ""
    return "||".join(parts)

def _get_tables_with_shapes(slide):
    """Return a list of (shape, table) pairs for all tables on the slide, including grouped shapes."""
    out = []
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shp = queue.pop(0)
        try:
            if hasattr(shp, "shapes"):
                try:
                    queue.extend(list(shp.shapes))
                except Exception:
                    pass
            if getattr(shp, "has_table", False):
                try:
                    out.append((shp, shp.table))
                except Exception:
                    pass
        except Exception:
            continue
    return out

def _collect_slide_text(slide) -> str:
    """Collect all text from a slide (including grouped shapes) into a single lowercase string."""
    texts = []
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shp = queue.pop(0)
        try:
            if hasattr(shp, "shapes"):
                try:
                    queue.extend(list(shp.shapes))
                except Exception:
                    pass
            if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                try:
                    texts.append(shp.text_frame.text or "")
                except Exception:
                    pass
        except Exception:
            continue
    try:
        joined = " ".join(" ".join(t.split()) for t in texts).strip().lower()
    except Exception:
        joined = ""
    return joined

def _find_slide_by_text_variants(prs, variants):
    """Find a slide whose aggregated text contains any of the provided variant strings (case-insensitive)."""
    wants = [" ".join(str(v or "").split()).strip().lower() for v in variants if v]
    for slide in prs.slides:
        txt = _collect_slide_text(slide)
        for w in wants:
            if w and w in txt:
                return slide
    return None

def _populate_entry_exit_table(table, df: pd.DataFrame, header_rows: int, mode: str) -> bool:
    """Populate an Entry/Exit table.

    mode: 'entry' or 'exit'
    """
    try:
        cat_col = _ci_col(df, ["Category", "category"]) or None
        if mode == "entry":
            name_col = _ci_col(df, ["New Shareholder", "New Share Holder", "Shareholder Name", "Name of Holder", "Name"]) or None
            shares_col = _ci_col(df, ["Shares Acquired during the Week", "Shares Acquired", "Shares Bought", "Buy Shares", "Shares"]) or None
        else:
            name_col = _ci_col(df, ["Exited Shareholder", "Exit Shareholder", "Shareholder Name", "Name of Holder", "Name"]) or None
            shares_col = _ci_col(df, ["Shares Sold during the Week", "Shares Sold", "Sold Shares", "Sell Shares", "Shares"]) or None
        pct_col = _ci_col(df, ["% of Share Capital", "% of Sh. Cap", "% of Sh. Cap (current)", "% of Sh. Cap (Current)", "Percentage"]) or None

        start_row = max(1, int(header_rows))
        _ensure_table_data_row_capacity(table, len(df), header_rows=header_rows)
        # Safety: enforce exact row count (header_rows + len(df)) by trimming any trailing rows
        try:
            total_needed = start_row + len(df)
            tbl_xml = table._tbl
            tr_list = list(tbl_xml.tr_lst)
            while len(tr_list) > total_needed:
                try:
                    tbl_xml.remove(tr_list[-1])
                    tr_list.pop()
                except Exception:
                    break
        except Exception:
            pass

        for r in range(start_row, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    table.cell(r, c).text = ""
                except Exception:
                    pass

        # Write exactly len(df) rows
        n = len(df)
        for i in range(n):
            row = df.iloc[i]
            r = start_row + i
            _write_data_cell(table, r, 0, (row[cat_col] if cat_col in df.columns else ""), align=PP_ALIGN.LEFT, font_size_pt=10, word_wrap=False)
            name_val = row[name_col] if (name_col in df.columns) else ""
            name_val = _normalize_shareholder_name(name_val)
            _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=10, allow_shrink=True, word_wrap=False)
            _write_data_cell(table, r, 2, (row[shares_col] if shares_col in df.columns else ""), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            _write_data_cell(table, r, 3, _format_pct(row[pct_col]) if pct_col in df.columns else "", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        _disable_table_autofit(table, start_row=0)
        _compact_data_rows(table, 0)
        _enforce_table_font_size(table, 0, font_size_pt=10)
        _shrink_column_fonts(table, 1, start_row=header_rows, base_pt=10, min_pt=6)
        # Post-trim: ensure no extra rows remain beyond header_rows + n
        try:
            total_needed = start_row + n
            tbl_xml = table._tbl
            while len(tbl_xml.tr_lst) > total_needed:
                tbl_xml.remove(tbl_xml.tr_lst[-1])
        except Exception:
            pass
        try:
            # Re-count from python-pptx object and underlying XML
            xml_rows = 0
            try:
                xml_rows = len(table._tbl.tr_lst)
            except Exception:
                pass
            print(f"[PPT][ENTRYEXIT] Wrote {n} data rows (start_row={start_row}, total_rows={len(table.rows)}, xml_rows={xml_rows})", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def update_new_entry_exits_on_template(prs: Presentation, df_entry: pd.DataFrame, df_exit: pd.DataFrame):
    """Update the 'New Entry / Exits' slide using Entry and Exit DB tables."""
    try:
        slide = _find_slide_by_title(prs, "New Entry / Exits")
        if slide is None:
            # Common title text variants
            title_variants = [
                "New Entry/ Exits",
                "New Entry/Exits",
                "New Entries / Exits",
                "New Entries & Exits",
                "New Entry & Exits",
                "New Entry and Exits",
                "Entry / Exit",
                "Entry & Exit",
            ]
            for v in title_variants:
                slide = _find_slide_by_title(prs, v)
                if slide is not None:
                    break
        if slide is None:
            # Fallback: search any text content on slides
            slide = _find_slide_by_text_variants(prs, [
                "new entry / exits",
                "new entry",
                "new entries",
            ])
        if slide is None:
            # Heuristic: a slide that has both words 'entry' and 'exit' somewhere and at least 2 tables
            best = None
            best_score = (-1, -1)  # (tables_count, text_score)
            for s in prs.slides:
                txt = _collect_slide_text(s)
                has_entry = ("entry" in txt or "entries" in txt)
                has_exit = ("exit" in txt or "exits" in txt)
                if has_entry and has_exit:
                    pairs = _get_tables_with_shapes(s)
                    score = (len(pairs), len(txt))
                    if score > best_score:
                        best = s
                        best_score = score
            slide = best
        if slide is None:
            return False

        pairs_all = _get_tables_with_shapes(slide)
        if not pairs_all:
            return False

        # Prefer matching by header keywords to avoid decorative tables
        entry_kw = ["new shareholder", "shares acquired", "% of share cap", "% of share capital"]
        exit_kw  = ["exited shareholder", "shares sold", "% of share cap", "% of share capital"]

        def _score_header(sig: str, kws):
            s = 0
            for k in kws:
                k2 = " ".join(k.split()).strip().lower()
                if k2 and k2 in sig:
                    s += 1
            return s

        ranked = []
        for shp, tbl in pairs_all:
            try:
                sig = _table_header_signature(tbl)
                e = _score_header(sig, entry_kw)
                x = _score_header(sig, exit_kw)
                w = int(getattr(shp, 'width', 0) or 0) * int(getattr(shp, 'height', 0) or 0)
                ranked.append((e, x, w, shp, tbl))
            except Exception:
                continue
        # Choose the best entry and exit tables by keyword match, breaking ties by area and left position
        entry_cands = sorted(ranked, key=lambda t: (t[0], t[2], -int(getattr(t[3], 'left', 0) or 0)), reverse=True)
        exit_cands  = sorted(ranked, key=lambda t: (t[1], t[2], -int(getattr(t[3], 'left', 0) or 0)), reverse=True)
        left_tbl = entry_cands[0][4] if entry_cands and entry_cands[0][0] > 0 else None
        right_tbl = exit_cands[0][4] if exit_cands and exit_cands[0][1] > 0 else None
        if left_tbl is None or right_tbl is None or left_tbl is right_tbl:
            # Fallback to two largest tables left-to-right
            pairs = sorted(pairs_all, key=lambda st: int(getattr(st[0], 'left', 0) or 0))
            if len(pairs) < 2:
                return False
            left_tbl = pairs[0][1]
            right_tbl = pairs[-1][1]
        if right_tbl is None or left_tbl is right_tbl:
            return False

        def _hdr_rows(_tbl):
            try:
                cols = len(_tbl.columns)
                if cols > 0 and len(_tbl.rows) > 1:
                    for c in range(min(6, cols)):
                        if (_tbl.cell(1, c).text or "").strip():
                            return 2
                return 1
            except Exception:
                return 1

        # The actual template slide has a single header row in each table.
        header_rows_left = 1
        header_rows_right = 1
        try:
            print(f"[PPT][ENTRYEXIT] Using left/right tables with header_rows=1", flush=True)
        except Exception:
            pass

        ok_left = _populate_entry_exit_table(left_tbl, df_entry if df_entry is not None else pd.DataFrame(), header_rows=header_rows_left, mode="entry")
        ok_right = _populate_entry_exit_table(right_tbl, df_exit if df_exit is not None else pd.DataFrame(), header_rows=header_rows_right, mode="exit")
        return bool(ok_left and ok_right)
    except Exception:
        return False

def _get_first_table(slide):
    """Return the first table object on a slide, if any."""
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shape = queue.pop(0)
        try:
            if getattr(shape, "has_table", False):
                return shape.table
            if hasattr(shape, "shapes"):
                try:
                    queue.extend(list(shape.shapes))
                except Exception:
                    pass
        except Exception:
            continue
    return None

def _get_best_table(slide):
    """Return the most likely data table on the slide: max columns and rows."""
    best = None
    best_score = (-1, -1)
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shape = queue.pop(0)
        try:
            if getattr(shape, "has_table", False):
                tbl = shape.table
                try:
                    area = int(getattr(shape, "width", 0)) * int(getattr(shape, "height", 0))
                except Exception:
                    area = 0
                score = (len(tbl.columns), len(tbl.rows), area)
                if score > best_score:
                    best = tbl
                    best_score = score
            if hasattr(shape, "shapes"):
                try:
                    queue.extend(list(shape.shapes))
                except Exception:
                    pass
        except Exception:
            continue
    return best

def _iter_tables(slide):
    """Yield all table objects on a slide, including inside group shapes."""
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shape = queue.pop(0)
        try:
            if getattr(shape, "has_table", False):
                yield shape.table
            if hasattr(shape, "shapes"):
                try:
                    queue.extend(list(shape.shapes))
                except Exception:
                    pass
        except Exception:
            continue

def _collect_header_tokens(table, header_rows: int = 2, max_cols: int = 10):
    tokens = set()
    try:
        hr = max(1, int(header_rows))
    except Exception:
        hr = 2
    try:
        mc = min(max_cols, len(table.columns))
    except Exception:
        mc = len(table.columns)
    try:
        rows_to_scan = min(hr, len(table.rows))
        for r in range(rows_to_scan):
            for c in range(mc):
                try:
                    txt = (table.cell(r, c).text or "").strip().lower()
                    if txt:
                        tokens.update(txt.replace("\n", " ").split())
                except Exception:
                    continue
    except Exception:
        pass
    return tokens

def _find_target_table(slide, mode: str):
    """Find the Buyers/Sellers data table by matching known header words.

    mode: 'buyers' or 'sellers'
    """
    mode = (mode or "").strip().lower()
    want_acq = (mode == "buyers")
    want_sold = (mode == "sellers")

    for tbl in _iter_tables(slide):
        # Build a concatenated header text from top 2 rows, first ~10 columns
        header_lines = []
        try:
            cols = len(tbl.columns)
        except Exception:
            cols = 0
        scan_cols = min(10, cols if cols else 0)
        try:
            total_rows = len(tbl.rows)
        except Exception:
            total_rows = 0
        scan_rows = min(2, total_rows)
        for r in range(scan_rows):
            parts = []
            for c in range(scan_cols):
                try:
                    parts.append((tbl.cell(r, c).text or "").strip())
                except Exception:
                    parts.append("")
            header_lines.append(" ".join([p for p in parts if p]))
        header_text = " ".join(header_lines).strip().lower()

        toks = _collect_header_tokens(tbl, header_rows=2, max_cols=10)

        # Shared expectations
        has_rank = "rank" in toks
        has_name = ("shareholder" in toks) or ("name" in toks)
        has_cat = "category" in toks
        # Either holding or % of share capital tokens usually appear
        has_hold = ("holding" in toks) or ("%" in header_text and "capital" in header_text)

        # Buyers vs Sellers specific header phrase
        buyers_phrase = "shares acquired during the week"
        sellers_phrase = "shares sold during the week"
        has_acq_phrase = buyers_phrase in header_text or "acquired" in header_text or "bought" in header_text
        has_sold_phrase = sellers_phrase in header_text or "sold" in header_text

        # Strong filters: 8 columns expected and at least ~22 rows (2 headers + 20 data)
        col_ok = (cols == 8)
        rows_ok = (total_rows >= 20)

        # If report dates are available, require both to appear in header text
        try:
            req_dates_ok = True
            if CURRENT_WEEK_DATE and PREVIOUS_WEEK_DATE:
                if isinstance(CURRENT_WEEK_DATE, str) and isinstance(PREVIOUS_WEEK_DATE, str):
                    req_dates_ok = (CURRENT_WEEK_DATE.strip().lower() in header_text) and (PREVIOUS_WEEK_DATE.strip().lower() in header_text)
        except Exception:
            req_dates_ok = True

        if has_rank and has_name and has_cat and has_hold and col_ok and rows_ok and req_dates_ok:
            if (want_acq and has_acq_phrase) or (want_sold and has_sold_phrase):
                return tbl
    # Fallback to best or first
    return _get_best_table(slide) or _get_first_table(slide)

def _set_cell_text_safe(table, r, c, value):
    if not (0 <= r < len(table.rows) and 0 <= c < len(table.columns)):
        return False
    try:
        table.cell(r, c).text = str(value if value is not None else "")
        return True
    except Exception:
        return False

def _fit_table_rows_to_shape_height(table, shape_height_emu: int, header_rows: int) -> bool:
    """Force table row heights to fit within a given shape height.

    This mimics manual drag-resize in PowerPoint (rows scale) which shape resizing alone
    may not reliably trigger in some templates.
    """
    try:
        total_h = int(shape_height_emu)
    except Exception:
        return False
    try:
        hr = max(0, int(header_rows))
    except Exception:
        hr = 0

    try:
        n_rows = len(table.rows)
        if n_rows <= 0:
            return False
        hr = min(hr, n_rows)
        data_rows = max(0, n_rows - hr)
        if data_rows <= 0:
            return True

        # Sum header heights (fallback to 0 if missing)
        header_sum = 0
        for r in range(hr):
            try:
                h = int(table.rows[r].height)
                if h > 0:
                    header_sum += h
            except Exception:
                continue

        avail = max(0, int(total_h) - int(header_sum))
        per = max(1, int(avail // data_rows))

        for r in range(hr, n_rows):
            try:
                target_h = int(per)
                table.rows[r].height = target_h
                try:
                    tr = table._tbl.tr_lst[r]
                    tr.set('h', str(int(target_h)))
                except Exception:
                    pass
            except Exception:
                continue
        try:
            print(f"[PPT] Fit data row heights to shape: header_sum_emu={header_sum}, per_data_row_emu={per}", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _disable_table_autofit(table, start_row: int = 0):
    """Disable PowerPoint text autofit for table cells so row heights can be enforced."""
    try:
        sr = max(0, int(start_row))
    except Exception:
        sr = 0
    try:
        for r in range(sr, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    cell = table.cell(r, c)
                    tf = getattr(cell, "text_frame", None)
                    if tf is None:
                        continue
                    try:
                        tf.auto_size = MSO_AUTO_SIZE.NONE
                    except Exception:
                        pass
                except Exception:
                    continue
        try:
            print(f"[PPT] Disabled autofit for table cells (rows={sr}..{len(table.rows)-1})", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _find_footer_legend_top(slide) -> Optional[int]:
    """Return top EMU of the footer legend textbox (FIIs/FPIs/DIIs...) if present."""
    try:
        queue = []
        try:
            queue = list(getattr(slide, "shapes", []))
        except Exception:
            queue = []
        best_top = None
        while queue:
            shp = queue.pop(0)
            try:
                if getattr(shp, "has_text_frame", False) and getattr(shp, "text_frame", None) is not None:
                    txt = (shp.text_frame.text or "").strip().lower()
                    if ("fiis:" in txt) or ("foreign institutional investors" in txt):
                        t = int(getattr(shp, "top", 0))
                        if best_top is None or t < best_top:
                            best_top = t
                if hasattr(shp, "shapes"):
                    try:
                        queue.extend(list(shp.shapes))
                    except Exception:
                        pass
            except Exception:
                continue
        return best_top
    except Exception:
        return None

def _set_table_height_to_avoid_footer(slide, table, top_inches: float, gap_inches: float = 0.06) -> bool:
    """Set table top (inches) and set height so bottom stays above the footer legend."""
    try:
        top_emu = Inches(float(top_inches))
        gap_emu = Inches(float(gap_inches))
    except Exception:
        return False

    footer_top = _find_footer_legend_top(slide)
    if footer_top is None:
        return False

    target_h_emu = max(0, int(footer_top) - int(top_emu) - int(gap_emu))
    try:
        _set_table_shape_top(slide, table, top_inches=top_inches)
    except Exception:
        pass
    try:
        _set_table_shape_height(slide, table, height_inches=float(target_h_emu) / 914400.0)
    except Exception:
        return False
    try:
        print(
            f"[PPT] Set table height to avoid footer: top_in={top_inches} footer_top_emu={int(footer_top)} -> height_in={target_h_emu/914400:.3f}",
            flush=True,
        )
    except Exception:
        pass
    return True

def _set_table_shape_bbox(slide, table, top_inches: float, bottom_margin_inches: float):
    """Set table shape top and height so bottom stays above slide footer area.

    Height is computed as: slide_height - top - bottom_margin.
    Also updates the parent group transform if the table is inside a group.
    """
    try:
        top_emu = Inches(float(top_inches))
        bottom_emu = Inches(float(bottom_margin_inches))
    except Exception:
        return False

    # Determine slide height (EMU)
    slide_h = None
    try:
        slide_h = int(getattr(getattr(slide, "part", None), "presentation", None).slide_height)
    except Exception:
        slide_h = None
    if not slide_h:
        try:
            # Fallback to common widescreen height
            slide_h = int(Inches(7.5))
        except Exception:
            return False

    target_h = max(0, int(slide_h) - int(top_emu) - int(bottom_emu))

    def _is_same_table(shp, tbl) -> bool:
        try:
            if not getattr(shp, "has_table", False):
                return False
            st = getattr(shp, "table", None)
            if st is None:
                return False
            return getattr(st, "_tbl", None) is getattr(tbl, "_tbl", None)
        except Exception:
            return False

    try:
        try:
            queue = [(shp, None) for shp in list(getattr(slide, "shapes", []))]
        except Exception:
            queue = []
        found = False
        while queue:
            shp, parent = queue.pop(0)
            try:
                if _is_same_table(shp, table):
                    # Apply on table shape
                    try:
                        shp.top = int(top_emu)
                    except Exception:
                        pass
                    try:
                        shp.height = int(target_h)
                    except Exception:
                        pass
                    try:
                        xfrm = getattr(getattr(shp, "_element", None), "xfrm", None)
                        if xfrm is not None:
                            if hasattr(xfrm, "off") and hasattr(xfrm.off, "y"):
                                xfrm.off.y = int(top_emu)
                            if hasattr(xfrm, "cy"):
                                xfrm.cy = int(target_h)
                    except Exception:
                        pass

                    # Apply on parent group if present
                    try:
                        if parent is not None:
                            try:
                                parent.top = int(top_emu)
                            except Exception:
                                pass
                            try:
                                parent.height = int(target_h)
                            except Exception:
                                pass
                            try:
                                px = getattr(getattr(parent, "_element", None), "xfrm", None)
                                if px is not None:
                                    if hasattr(px, "off") and hasattr(px.off, "y"):
                                        px.off.y = int(top_emu)
                                    if hasattr(px, "cy"):
                                        px.cy = int(target_h)
                            except Exception:
                                pass
                    except Exception:
                        pass

                    try:
                        print(
                            f"[PPT] Set table bbox top={top_inches}in bottom_margin={bottom_margin_inches}in -> height_in={target_h/914400:.3f}",
                            flush=True,
                        )
                    except Exception:
                        pass
                    found = True
                    break

                if hasattr(shp, "shapes"):
                    try:
                        queue.extend([(ch, shp) for ch in list(shp.shapes)])
                    except Exception:
                        pass
            except Exception:
                continue
        if not found:
            try:
                print("[PPT][WARN] Table shape not found for setting bbox", flush=True)
            except Exception:
                pass
        return True
    except Exception:
        return False

def _nudge_table_shape_up(slide, table, delta_inches: float):
    """Move the table shape up by delta_inches (and its parent group if grouped)."""
    try:
        delta = Inches(float(delta_inches))
    except Exception:
        return False

    def _is_same_table(shp, tbl) -> bool:
        try:
            if not getattr(shp, "has_table", False):
                return False
            st = getattr(shp, "table", None)
            if st is None:
                return False
            return getattr(st, "_tbl", None) is getattr(tbl, "_tbl", None)
        except Exception:
            return False

    try:
        try:
            queue = [(shp, None) for shp in list(getattr(slide, "shapes", []))]
        except Exception:
            queue = []
        found = False
        while queue:
            shp, parent = queue.pop(0)
            try:
                if _is_same_table(shp, table):
                    try:
                        new_top = max(0, int(shp.top) - int(delta))
                    except Exception:
                        new_top = None
                    if new_top is not None:
                        try:
                            shp.top = new_top
                        except Exception:
                            pass
                        try:
                            xfrm = getattr(getattr(shp, "_element", None), "xfrm", None)
                            if xfrm is not None and hasattr(xfrm, "off") and hasattr(xfrm.off, "y"):
                                xfrm.off.y = int(new_top)
                        except Exception:
                            pass
                        try:
                            if parent is not None:
                                parent.top = new_top
                                try:
                                    px = getattr(getattr(parent, "_element", None), "xfrm", None)
                                    if px is not None and hasattr(px, "off") and hasattr(px.off, "y"):
                                        px.off.y = int(new_top)
                                except Exception:
                                    pass
                        except Exception:
                            pass
                        try:
                            print(f"[PPT] Nudged table up by {delta_inches} in (new_top_emu={int(new_top)})", flush=True)
                        except Exception:
                            pass
                    found = True
                    break
                if hasattr(shp, "shapes"):
                    try:
                        queue.extend([(ch, shp) for ch in list(shp.shapes)])
                    except Exception:
                        pass
            except Exception:
                continue
        if not found:
            try:
                print("[PPT][WARN] Table shape not found for nudging up", flush=True)
            except Exception:
                pass
        return True
    except Exception:
        return False

def _enable_table_text_to_fit(table, start_row: int = 0):
    """Enable 'Text to fit shape' for table cells so content shrinks into a fixed table size."""
    try:
        sr = max(0, int(start_row))
    except Exception:
        sr = 0
    try:
        for r in range(sr, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    cell = table.cell(r, c)
                    tf = getattr(cell, "text_frame", None)
                    if tf is None:
                        continue
                    try:
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass
                except Exception:
                    continue
        try:
            print(f"[PPT] Enabled text-to-fit for table cells (rows={sr}..{len(table.rows)-1})", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _shrink_column_fonts(table, col_index: int, start_row: int, base_pt: int = 10, min_pt: int = 6):
    """Shrink font size for a specific table column starting at start_row to reduce wrapping."""
    try:
        sr = max(0, int(start_row))
    except Exception:
        sr = 0
    try:
        ci = int(col_index)
    except Exception:
        return False
    if ci < 0:
        return False
    try:
        for r in range(sr, len(table.rows)):
            if ci >= len(table.columns):
                break
            try:
                cell = table.cell(r, ci)
                tf = getattr(cell, "text_frame", None)
                if tf is None:
                    continue
                # Keep no-wrap to avoid multi-line forcing row expansion
                try:
                    tf.word_wrap = False
                except Exception:
                    pass
                txt = ""
                try:
                    txt = " ".join((tf.text or "").split())
                except Exception:
                    txt = str(tf.text or "")
                ln = len(txt)
                size = int(base_pt)
                # More aggressive than default: target single-line names
                if ln > 70:
                    size = max(min_pt, base_pt - 5)
                elif ln > 60:
                    size = max(min_pt, base_pt - 4)
                elif ln > 50:
                    size = max(min_pt, base_pt - 3)
                elif ln > 42:
                    size = max(min_pt, base_pt - 2)
                elif ln > 34:
                    size = max(min_pt, base_pt - 1)
                for para in getattr(tf, "paragraphs", []):
                    try:
                        para.font.size = Pt(size)
                    except Exception:
                        pass
                    for run in getattr(para, "runs", []):
                        try:
                            run.font.size = Pt(size)
                            run.font.name = "Arial"
                        except Exception:
                            pass
                # Update underlying XML so selection reflects size
                try:
                    tx_body = tf._txBody
                    for elem in tx_body.iter():
                        if elem.tag in (qn('a:rPr'), qn('a:defRPr'), qn('a:endParaRPr')):
                            elem.set('sz', str(int(size * 100)))
                except Exception:
                    pass
            except Exception:
                continue
        try:
            print(f"[PPT] Shrunk column {ci} fonts from row {sr} (base={base_pt}, min={min_pt})", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _enforce_table_font_size(table, start_row: int, font_size_pt: int = 10):
    """Force all runs in data rows of a table to a fixed font size (e.g., 10pt)."""
    try:
        sr = max(0, int(start_row))
    except Exception:
        sr = 0
    try:
        for r in range(sr, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    cell = table.cell(r, c)
                    tf = getattr(cell, "text_frame", None)
                    if tf is None:
                        continue
                    for para in getattr(tf, "paragraphs", []):
                        # Set paragraph-level font size so selection reflects 10pt consistently
                        try:
                            para.font.size = Pt(font_size_pt)
                        except Exception:
                            pass
                        runs = list(getattr(para, "runs", []))
                        # Ensure at least one run exists so PowerPoint reflects uniform size
                        if not runs:
                            try:
                                new_run = para.add_run()
                                new_run.text = ""
                                new_run.font.size = Pt(font_size_pt)
                                new_run.font.name = "Arial"
                            except Exception:
                                pass
                            runs = list(getattr(para, "runs", []))
                        for run in runs:
                            try:
                                run.font.size = Pt(font_size_pt)
                                run.font.name = "Arial"
                            except Exception:
                                pass

                    # Also update underlying XML run properties so whole-table selection shows this size
                    try:
                        tx_body = tf._txBody
                        for elem in tx_body.iter():
                            if elem.tag in (qn('a:rPr'), qn('a:defRPr'), qn('a:endParaRPr')):
                                elem.set('sz', str(int(font_size_pt * 100)))
                    except Exception:
                        pass
                except Exception:
                    continue
        return True
    except Exception:
        return False

def _enforce_column_font_size(table, col_index: int, start_row: int, font_size_pt: int = 9):
    """Force a specific table column to a fixed font size, including XML run properties."""
    try:
        sr = max(0, int(start_row))
    except Exception:
        sr = 0
    try:
        ci = int(col_index)
    except Exception:
        return False
    if ci < 0:
        return False
    try:
        for r in range(sr, len(table.rows)):
            if ci >= len(table.columns):
                break
            try:
                cell = table.cell(r, ci)
                tf = getattr(cell, "text_frame", None)
                if tf is None:
                    continue
                try:
                    tf.auto_size = MSO_AUTO_SIZE.NONE
                except Exception:
                    pass
                try:
                    tf.word_wrap = False
                except Exception:
                    pass
                for para in getattr(tf, "paragraphs", []):
                    try:
                        para.font.size = Pt(font_size_pt)
                    except Exception:
                        pass
                    runs = list(getattr(para, "runs", []))
                    if not runs:
                        try:
                            new_run = para.add_run()
                            new_run.text = ""
                            new_run.font.size = Pt(font_size_pt)
                            new_run.font.name = "Arial"
                        except Exception:
                            pass
                        runs = list(getattr(para, "runs", []))
                    for run in runs:
                        try:
                            run.font.size = Pt(font_size_pt)
                            run.font.name = "Arial"
                        except Exception:
                            pass
                try:
                    tx_body = tf._txBody
                    for elem in tx_body.iter():
                        if elem.tag in (qn('a:rPr'), qn('a:defRPr'), qn('a:endParaRPr')):
                            elem.set('sz', str(int(font_size_pt * 100)))
                    body_pr = getattr(tx_body, 'bodyPr', None)
                    if body_pr is not None:
                        try:
                            body_pr.set('wrap', 'none')
                        except Exception:
                            pass
                    norm_text = " ".join((tf.text or "").split())
                    if norm_text != (tf.text or ""):
                        try:
                            tf.clear()
                            p = tf.paragraphs[0]
                            p.alignment = PP_ALIGN.LEFT
                            rr = p.add_run()
                            rr.text = norm_text
                            rr.font.size = Pt(font_size_pt)
                            rr.font.name = "Arial"
                        except Exception:
                            pass
                except Exception:
                    pass
            except Exception:
                continue
        return True
    except Exception:
        return False

def _set_header_row_height(table, header_rows: int, height_inches: float):
    """Set a uniform height (in inches) for the header rows (0..header_rows-1)."""
    try:
        hr = max(0, int(header_rows))
    except Exception:
        hr = 0
    if hr <= 0:
        return True
    try:
        try:
            target_h = Inches(height_inches)
        except Exception:
            return False
        rows_to_set = min(hr, len(table.rows))
        for r in range(rows_to_set):
            try:
                row = table.rows[r]
                row.height = target_h
                try:
                    tr = table._tbl.tr_lst[r]
                    tr.set('h', str(int(target_h)))
                except Exception:
                    pass
            except Exception:
                continue
        try:
            print(f"[PPT] Header row heights set: rows=0..{rows_to_set-1}, h_in={height_inches}", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _set_table_shape_top(slide, table, top_inches: float):
    """Position the table's shape from the top of the slide (in inches)."""
    try:
        target_t = Inches(top_inches)
    except Exception:
        return False

    def _is_same_table(shp, tbl) -> bool:
        try:
            if not getattr(shp, "has_table", False):
                return False
            st = getattr(shp, "table", None)
            if st is None:
                return False
            # Compare underlying XML table nodes; object identity may differ
            return getattr(st, "_tbl", None) is getattr(tbl, "_tbl", None)
        except Exception:
            return False

    try:
        queue = []
        try:
            queue = list(getattr(slide, "shapes", []))
        except Exception:
            queue = []
        found = False
        while queue:
            shp = queue.pop(0)
            try:
                if _is_same_table(shp, table):
                    shp.top = target_t
                    # Defensive: also set underlying transform offset if available
                    try:
                        xfrm = getattr(getattr(shp, "_element", None), "xfrm", None)
                        if xfrm is not None and hasattr(xfrm, "off") and hasattr(xfrm.off, "y"):
                            xfrm.off.y = int(target_t)
                    except Exception:
                        pass
                    try:
                        print(f"[PPT] Set table top to {top_inches} in (EMU={int(target_t)})", flush=True)
                    except Exception:
                        pass
                    found = True
                    break
                if hasattr(shp, "shapes"):
                    try:
                        queue.extend(list(shp.shapes))
                    except Exception:
                        pass
            except Exception:
                continue
        if not found:
            try:
                print("[PPT][WARN] Table shape not found for setting top", flush=True)
            except Exception:
                pass
        return True
    except Exception:
        return False

def _compact_data_rows(table, header_rows: int):
    try:
        hr = max(0, int(header_rows))
    except Exception:
        hr = 0
    try:
        for r in range(hr, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    cell = table.cell(r, c)
                    try:
                        cell.margin_left = Emu(0)
                        cell.margin_right = Emu(0)
                        cell.margin_top = Emu(0)
                        cell.margin_bottom = Emu(0)
                    except Exception:
                        pass
                    tf = getattr(cell, "text_frame", None)
                    if tf is None:
                        continue
                    for para in getattr(tf, "paragraphs", []):
                        try:
                            para.space_before = Pt(0)
                            para.space_after = Pt(0)
                        except Exception:
                            pass
                except Exception:
                    continue
        return True
    except Exception:
        return False

def _set_data_row_height(table, header_rows: int, height_inches: float):
    """Set a uniform height (in inches) for all data rows below header_rows.

    This is used to compress data rows on dense tables (like Top 20 Buyers/Sellers)
    so that all rows fit comfortably within the slide while preserving headers.
    """
    try:
        hr = max(0, int(header_rows))
    except Exception:
        hr = 0
    try:
        target_h = Inches(height_inches)
    except Exception:
        return False
    try:
        for r in range(hr, len(table.rows)):
            try:
                row = table.rows[r]
                # Set python-pptx row height
                row.height = target_h
                # Also force the underlying XML row height attribute
                try:
                    tr = table._tbl.tr_lst[r]
                    tr.set('h', str(int(target_h)))
                except Exception:
                    pass
            except Exception:
                continue
        # Defensive: force the XML tr heights for all data rows even if row objects are quirky
        try:
            tr_lst = getattr(table._tbl, "tr_lst", [])
            for r in range(hr, min(len(tr_lst), len(table.rows))):
                try:
                    tr_lst[r].set('h', str(int(target_h)))
                except Exception:
                    pass
        except Exception:
            pass
        try:
            print(f"[PPT] Data row heights set: rows={hr}..{len(table.rows)-1}, h_in={height_inches}", flush=True)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _set_table_shape_height(slide, table, height_inches: float):
    """Set the overall height of the table's shape so rows don't re-expand.

    PowerPoint often stretches rows to fill the table shape height. By reducing
    the shape height we can compress row spacing visually while preserving the
    template layout.
    """
    try:
        target_h = Inches(height_inches)
    except Exception:
        return False

    def _is_same_table(shp, tbl) -> bool:
        try:
            if not getattr(shp, "has_table", False):
                return False
            st = getattr(shp, "table", None)
            if st is None:
                return False
            return getattr(st, "_tbl", None) is getattr(tbl, "_tbl", None)
        except Exception:
            return False

    try:
        # Search recursively in case the table is inside a group shape (keep parent so we can resize it too)
        queue = []
        try:
            queue = [(shp, None) for shp in list(getattr(slide, "shapes", []))]
        except Exception:
            queue = []
        found = False
        while queue:
            shp, parent = queue.pop(0)
            try:
                if _is_same_table(shp, table):
                    shp.height = target_h
                    # Defensive: also set underlying transform cy if available
                    try:
                        xfrm = getattr(getattr(shp, "_element", None), "xfrm", None)
                        if xfrm is not None and hasattr(xfrm, "cy"):
                            xfrm.cy = int(target_h)
                    except Exception:
                        pass
                    # If the table is inside a group, also resize the group height so PPT doesn't clamp it
                    try:
                        if parent is not None:
                            parent.height = target_h
                            try:
                                px = getattr(getattr(parent, "_element", None), "xfrm", None)
                                if px is not None and hasattr(px, "cy"):
                                    px.cy = int(target_h)
                            except Exception:
                                pass
                    except Exception:
                        pass
                    try:
                        print(f"[PPT] Set table height to {height_inches} in (EMU={int(target_h)})", flush=True)
                    except Exception:
                        pass
                    found = True
                    break
                if hasattr(shp, "shapes"):
                    try:
                        queue.extend([(ch, shp) for ch in list(shp.shapes)])
                    except Exception:
                        pass
            except Exception:
                continue
        if not found:
            try:
                print("[PPT][WARN] Table shape not found for setting height", flush=True)
            except Exception:
                pass
        return True
    except Exception:
        return False

def update_insurance_pf_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 10 DII - Insurance & PFs' table on the template slide (page 10).
    Targets slide index 9 first, then falls back to title search.
    Column mapping per user:
      - Rank -> serial.no (also accept Rank)
      - Shareholder Name -> Institution
      - Use two latest date-named columns for Previous (e.g., 12/19/2025) and Current (e.g., 12/26/2025)
      - % of Share Capital -> % of Sh. Cap (Previous)/(Current)
      - Change in holding shares -> MoM change in holdings
    """
    # Prefer slide 10 (0-based index 9)
    slide = None
    try:
        if len(prs.slides) > 9:
            slide = prs.slides[9]
    except Exception:
        slide = None
    if slide is None:
        # Try various common title variants
        candidates = [
            "Top 10 DII - Insurance & PFs",
            "Top 10 Insurance & PFs",
            "Top 10 DII - PFs",
            "Top 10 DII â€“ Insurance & PFs",
        ]
        for title in candidates:
            sld = _find_slide_by_title(prs, title)
            if sld is not None:
                slide = sld
                break
    if slide is None:
        # Any slide whose TITLE placeholder or any text contains both 'insurance'/'pfs'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            txt = (" ".join(shp.text_frame.text.split()).strip().lower())
                            if ("insurance" in txt and ("pf" in txt or "pfs" in txt)):
                                slide = s
                                break
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame and slide is None:
                        raw = (shp.text_frame.text or "").strip().lower()
                        txt = " ".join(raw.split())
                        if ("insurance" in txt and ("pf" in txt or "pfs" in txt)):
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        return False

    # Announce which slide
    try:
        slide_idx = None
        for i, s in enumerate(prs.slides):
            if s == slide:
                slide_idx = i
                break
        if slide_idx is not None:
            print(f"   â„¹ï¸ Updating 'Top 10 DII - Insurance & PFs' on slide #{slide_idx + 1}")
    except Exception:
        pass

    table = _get_best_table(slide) or _get_first_table(slide)
    if table is None:
        return False
    if len(table.columns) < 7:
        return False

    # Column mappings
    rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Institution", "Shareholder Name"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None
    buy_pps_col  = _ci_col(df, ["Buy (pps)", "Buy pps", "Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell (pps)", "Sell pps", "Sell"]) or None
    # Optional buy/sell pps columns for coloring Change
    buy_pps_col  = _ci_col(df, ["Buy (pps)", "Buy pps", "Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell (pps)", "Sell pps", "Sell"]) or None

    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(6, len(table.columns)))) else 1
    left_is_prev = False if _force_current_left() else True
    if not _force_current_left():
        try:
            hdr_left_txt = (table.cell(0, 2).text or "").strip() if len(table.columns) > 2 else ""
            hdr_right_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
            dl = _parse_date_label_to_dt(hdr_left_txt)
            dr = _parse_date_label_to_dt(hdr_right_txt)
            if dl and dr:
                left_is_prev = dl < dr
        except Exception:
            pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 2, 4, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 2, 4, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    # Clear data rows
    start_row = header_rows
    _ensure_table_data_row_capacity(table, len(df), header_rows=start_row)
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0
    for i in range(n):
        row = df.iloc[i]
        r = start_row + i

        # Rank
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Shareholder Name
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)

        if left_is_prev:
            prev_block_cols = (2, 3)
            curr_block_cols = (4, 5)
        else:
            curr_block_cols = (2, 3)
            prev_block_cols = (4, 5)

        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Change in Holding (pps): prefer Buy (green), else Sell (red), else fallback
        def _nonzero(x):
            try:
                s = str(x).strip()
                if not s or s.lower() == 'none':
                    return False
                v = float(s.replace(',', '').replace('(', '-').replace(')', ''))
                return abs(v) > 0
            except Exception:
                return bool(str(x).strip())
        buy_v = (row[buy_pps_col] if (buy_pps_col in df.columns) else None) if buy_pps_col is not None else None
        sell_v = (row[sell_pps_col] if (sell_pps_col in df.columns) else None) if sell_pps_col is not None else None
        if _nonzero(buy_v):
            _write_data_cell(table, r, 6, _format_change_value(buy_v), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        elif _nonzero(sell_v):
            _write_data_cell(table, r, 6, f"({_format_change_value(sell_v)})", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=RED_COLOR)
        else:
            _write_data_cell(table, r, 6, "-", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    return True

def update_aif_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 10 DII - AIFs' table on the template slide (page 11).
    Targets slide index 10 first, then falls back to title search.
    Column mapping per user:
      - Rank -> serial.no (also accept Rank)
      - Shareholder Name -> Institution
      - Use two latest date-named columns for Previous and Current
      - % of Share Capital -> % of Sh. Cap (Previous)/(Current)
      - Change in holding shares -> MoM change in holdings
    """
    # Prefer slide 11 (0-based index 10)
    slide = None
    try:
        if len(prs.slides) > 10:
            slide = prs.slides[10]
    except Exception:
        slide = None
    if slide is None:
        # Try common title variants
        candidates = [
            "Top 10 DII - AIFs",
            "Top 10 AIFs",
            "Top 10 DII â€“ AIFs",
        ]
        for title in candidates:
            sld = _find_slide_by_title(prs, title)
            if sld is not None:
                slide = sld
                break
    if slide is None:
        # Any slide whose TITLE placeholder or any text contains 'aif'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            txt = (" ".join(shp.text_frame.text.split()).strip().lower())
                            if ("aif" in txt):
                                slide = s
                                break
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame and slide is None:
                        raw = (shp.text_frame.text or "").strip().lower()
                        txt = " ".join(raw.split())
                        if ("aif" in txt):
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        return False

    # Announce which slide
    try:
        slide_idx = None
        for i, s in enumerate(prs.slides):
            if s == slide:
                slide_idx = i
                break
        if slide_idx is not None:
            print(f"   â„¹ï¸ Updating 'Top 10 DII - AIFs' on slide #{slide_idx + 1}")
    except Exception:
        pass

    table = _get_best_table(slide) or _get_first_table(slide)
    if table is None:
        return False
    if len(table.columns) < 7:
        return False

    # Column mappings
    rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Institution", "Shareholder Name"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None

    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(6, len(table.columns)))) else 1
    left_is_prev = False if _force_current_left() else True
    if not _force_current_left():
        try:
            hdr_left_txt = (table.cell(0, 2).text or "").strip() if len(table.columns) > 2 else ""
            hdr_right_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
            dl = _parse_date_label_to_dt(hdr_left_txt)
            dr = _parse_date_label_to_dt(hdr_right_txt)
            if dl and dr:
                left_is_prev = dl < dr
        except Exception:
            pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 2, 4, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 2, 4, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    # Clear data rows
    start_row = header_rows
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0
    for i in range(n):
        row = df.iloc[i]
        r = start_row + i

        # Rank
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Shareholder Name
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)

        if left_is_prev:
            prev_block_cols = (2, 3)
            curr_block_cols = (4, 5)
        else:
            curr_block_cols = (2, 3)
            prev_block_cols = (4, 5)

        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        def _nonzero(x):
            try:
                s = str(x).strip()
                if not s or s.lower() == 'none':
                    return False
                v = float(s.replace(',', '').replace('(', '-').replace(')', ''))
                return abs(v) > 0
            except Exception:
                return bool(str(x).strip())
        buy_v = (row[buy_pps_col] if (buy_pps_col in df.columns) else None) if buy_pps_col is not None else None
        sell_v = (row[sell_pps_col] if (sell_pps_col in df.columns) else None) if sell_pps_col is not None else None
        if _nonzero(buy_v):
            _write_data_cell(table, r, 6, _format_change_value(buy_v), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        elif _nonzero(sell_v):
            _write_data_cell(table, r, 6, f"({_format_change_value(sell_v)})", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=RED_COLOR)
        else:
            _write_data_cell(table, r, 6, "-", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    return True

def update_mf_tables_on_template(prs: Presentation, df_active: pd.DataFrame, df_passive: pd.DataFrame):
    # Prefer explicit slide index 9 (0-based 8) as requested
    slide = None
    try:
        if len(prs.slides) > 8:
            slide = prs.slides[8]
    except Exception:
        slide = None
    if slide is None:
        slide = _find_slide_by_title(prs, "Top 10 DII - MFs")
    if slide is None:
        slide = _find_slide_by_title(prs, "Top 10 DII-MFs") or _find_slide_by_title(prs, "Top 10 MFs") or _find_slide_by_title(prs, "Top 10 MF's")
    if slide is None:
        # Fallback A: any slide whose TITLE placeholder contains 'mf'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            txt = (" ".join(shp.text_frame.text.split()).strip().lower())
                            if "mf" in txt:
                                slide = s
                                break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        # Fallback B: scan any text frame for the exact title or DII/MF keywords
        want = "top 10 dii - mfs"
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        raw = (shp.text_frame.text or "").strip().lower()
                        txt = " ".join(raw.split())
                        txt = txt.replace("â€“", "-").replace("â€”", "-").replace("â€™", "'")
                        if txt == want or ("dii" in txt and "mf" in txt):
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        try:
            print("   âš ï¸  Could not find 'Top 10 DII - MFs' slide by title or text search")
        except Exception:
            pass
        return False

    # Log which slide index (1-based) is being updated for clarity
    try:
        slide_idx = None
        for i, s in enumerate(prs.slides):
            if s == slide:
                slide_idx = i
                break
        if slide_idx is not None:
            print(f"   â„¹ï¸ Updating 'Top 10 DII - MFs' on slide #{slide_idx + 1}")
    except Exception:
        pass

    tables = []
    queue = []
    try:
        queue = list(getattr(slide, "shapes", []))
    except Exception:
        queue = []
    while queue:
        shape = queue.pop(0)
        try:
            if getattr(shape, "has_table", False):
                tbl = shape.table
                try:
                    area = int(getattr(shape, "width", 0)) * int(getattr(shape, "height", 0))
                except Exception:
                    area = 0
                top = int(getattr(shape, "top", 0)) if hasattr(shape, "top") else 0
                if len(tbl.columns) >= 7:
                    tables.append((tbl, area, top))
            if hasattr(shape, "shapes"):
                try:
                    queue.extend(list(shape.shapes))
                except Exception:
                    pass
        except Exception:
            continue
    if len(tables) == 0:
        try:
            print("   âš ï¸  No tables detected on 'Top 10 DII - MFs' slide")
        except Exception:
            pass
        return False
    tables.sort(key=lambda x: (x[2], -x[1]))
    table_active = tables[0][0]
    table_passive = tables[1][0] if len(tables) >= 2 else None

    def _update_7col_table(table, df: pd.DataFrame, squeeze: bool = False):
        if table is None or df is None or df.empty:
            return False
        if len(table.columns) < 7:
            return False
        header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(6, len(table.columns)))) else 1
        try:
            df = df.head(10).copy()
        except Exception:
            pass
        try:
            _ensure_table_data_row_capacity(table, len(df), header_rows=header_rows)
        except Exception:
            pass
        left_is_prev = False if _force_current_left() else True
        if not _force_current_left():
            try:
                hdr_left_txt = (table.cell(0, 2).text or "").strip() if len(table.columns) > 2 else ""
                hdr_right_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
                dl = _parse_date_label_to_dt(hdr_left_txt)
                dr = _parse_date_label_to_dt(hdr_right_txt)
                if dl and dr:
                    left_is_prev = dl < dr
            except Exception:
                pass
        for r in range(header_rows, len(table.rows)):
            for c in range(len(table.columns)):
                try:
                    table.cell(r, c).text = ""
                except Exception:
                    pass
        rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
        name_col = _ci_col(df, ["Institution", "Shareholder Name", "Short Name", "Name of Holder"]) or None
        prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
        prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
        curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])
        buy_pps_col  = _ci_col(df, ["Buy"]) or None
        sell_pps_col = _ci_col(df, ["Sell"]) or None
        try:
            if left_is_prev:
                _apply_header_date_labels(table, 2, 4, prev_hold_col, curr_hold_col)
            else:
                _apply_header_date_labels(table, 2, 4, curr_hold_col, prev_hold_col)
        except Exception:
            pass
        start_row = header_rows
        max_rows = len(table.rows) - start_row
        try:
            df = _sort_df_rankwise(df)
        except Exception:
            pass
        n = min(max_rows, len(df)) if max_rows > 0 else 0
        for i in range(n):
            row = df.iloc[i]
            r = start_row + i
            val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
            _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            name_val = row[name_col] if name_col in df.columns else ""
            name_val = _normalize_shareholder_name(name_val)
            _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)
            if left_is_prev:
                prev_block_cols = (2, 3)
                curr_block_cols = (4, 5)
            else:
                curr_block_cols = (2, 3)
                prev_block_cols = (4, 5)
            prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
            prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
            _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
            curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
            _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            # Change in Holding (pps) from Buy/Sell with colors; fallback to change_col
            def _nonzero(x):
                try:
                    s = str(x).strip()
                    if not s or s.lower() == 'none':
                        return False
                    v = float(s.replace(',', '').replace('(', '-').replace(')', ''))
                    return abs(v) > 0
                except Exception:
                    return bool(str(x).strip())
            buy_v = (row[buy_pps_col] if (buy_pps_col in df.columns) else None) if buy_pps_col is not None else None
            sell_v = (row[sell_pps_col] if (sell_pps_col in df.columns) else None) if sell_pps_col is not None else None
            if _nonzero(buy_v):
                _write_data_cell(table, r, 6, _format_change_value(buy_v), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            elif _nonzero(sell_v):
                _write_data_cell(table, r, 6, f"({_format_change_value(sell_v)})", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=RED_COLOR)
            else:
                _write_data_cell(table, r, 6, "-", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        try:
            _compact_data_rows(table, header_rows=header_rows)
        except Exception:
            pass
        if squeeze:
            try:
                _set_data_row_height(table, header_rows=header_rows, height_inches=0.20)
            except Exception:
                pass
            try:
                _set_table_shape_height(slide, table, height_inches=2.55)
            except Exception:
                pass
        return True

    ok1 = _update_7col_table(table_active, df_active, squeeze=False)
    ok2 = _update_7col_table(table_passive, df_passive, squeeze=True) if table_passive is not None else False
    return bool(ok1 or ok2)

def update_fii_fpi_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 10 FIIs & FPIs' table on the template slide with DB data (update-only).
    Mappings provided by user:
      - Rank -> serial.no (also accept Rank)
      - Shareholder Name -> Institution (also accept Shareholder Name / Short Name)
      - Current/Previous blocks -> two latest date-named columns in DF
      - % of Share Capital -> map to '% of Sh. Cap (Current)' and '(Previous)'
      - Change in holding shares -> 'MoM change in holdings'
    """
    try:
        if df is None:
            df = pd.DataFrame()
        else:
            df = df.head(10).copy()
    except Exception:
        pass

    # 1) Prefer explicit slide index 8 (0-based 7) per user instruction
    slide = None
    try:
        if len(prs.slides) > 7:
            slide = prs.slides[7]
    except Exception:
        slide = None

    # If not available, locate the FII/FPI slide by likely titles/keywords
    candidates = [
        "Top 10 FII's & FPI's",
        "Top 10 FIIs & FPIs",
        "Top 20 FII's & FPI's",
        "Top 20 FIIs & FPIs",
    ]
    if slide is None:
        for title in candidates:
            sld = _find_slide_by_title(prs, title)
            if sld is not None:
                slide = sld
                break
    if slide is None:
        # Fallback A: any slide whose TITLE placeholder contains both 'fii' and 'fpi'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            txt = (" ".join(shp.text_frame.text.split()).strip().lower())
                            if ("fii" in txt) and ("fpi" in txt):
                                slide = s
                                break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        # Fallback B: scan any text frame
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = (shp.text_frame.text or "").strip().lower()
                        if ("fii" in txt) and ("fpi" in txt):
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        return False

    # 2) Identify the main data table (expect 7 columns)
    table = _get_best_table(slide) or _get_first_table(slide)
    if table is None:
        return False
    if len(table.columns) < 7:
        return False

    # 3) Resolve DF column mappings (case-insensitive exact first)
    rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Institution", "Shareholder Name", "Short Name"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None

    # 4) Determine if the LEFT date block is previous or current by parsing header dates at top row
    #    For 7-col table, header row 0 has date labels at col 2 (current) and col 4 (previous) in our template
    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(6, len(table.columns)))) else 1
    left_is_prev = False if _force_current_left() else True
    if not _force_current_left():
        try:
            hdr_left_txt = (table.cell(0, 2).text or "").strip() if len(table.columns) > 2 else ""
            hdr_right_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
            dl = _parse_date_label_to_dt(hdr_left_txt)
            dr = _parse_date_label_to_dt(hdr_right_txt)
            if dl and dr:
                left_is_prev = dl < dr
        except Exception:
            pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 2, 4, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 2, 4, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    # 5) Clear existing data rows and write from DF
    # Ensure table has exactly header_rows + len(df) rows so extra blank rows don't stretch height
    _ensure_table_data_row_capacity(table, len(df), header_rows=header_rows)
    start_row = header_rows
    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    for i in range(n):
        row = df.iloc[i]
        r = start_row + i
        # Col 0: Rank
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        # Col 1: Shareholder Name (normalized)
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)

        # Determine date block positions
        if left_is_prev:
            prev_block_cols = (2, 3)
            curr_block_cols = (4, 5)
        else:
            curr_block_cols = (2, 3)
            prev_block_cols = (4, 5)

        # Previous block
        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Current block
        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Change column (col 6): from Buy/Sell with color; fallback to change_col
        def _nonzero(x):
            try:
                s = str(x).strip()
                if not s or s.lower() == 'none':
                    return False
                v = float(s.replace(',', '').replace('(', '-').replace(')', ''))
                return abs(v) > 0
            except Exception:
                return bool(str(x).strip())
        buy_v = (row[buy_pps_col] if (buy_pps_col in df.columns) else None) if buy_pps_col is not None else None
        sell_v = (row[sell_pps_col] if (sell_pps_col in df.columns) else None) if sell_pps_col is not None else None
        if _nonzero(buy_v):
            _write_data_cell(table, r, 6, _format_change_value(buy_v), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        elif _nonzero(sell_v):
            _write_data_cell(table, r, 6, f"({_format_change_value(sell_v)})", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=RED_COLOR)
        else:
            _write_data_cell(table, r, 6, "-", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    return True

def update_buyers_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 20 Buyers' table on the template slide with DB data (update-only)."""
    try:
        print("[PPT][TRACE] Enter update_buyers_table_on_template", flush=True)
    except Exception:
        pass
    # Locate the slide by exact title first, then fallbacks
    slide = _find_slide_by_title(prs, "Top 20 Buyers")
    if slide is None:
        # Fallback A: any slide whose TITLE placeholder contains 'buyers'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            if "buyers" in (" ".join(shp.text_frame.text.split()).strip().lower()):
                                slide = s
                                break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        # Fallback B: scan any text frame for 'top 20 buyers'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = (shp.text_frame.text or "").strip().lower()
                        if "top 20 buyers" in txt or "buyers" in txt:
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None and len(prs.slides) > 2:
        # Conservative default: third slide if template is known order
        slide = prs.slides[2]
    if slide is None:
        return False

    table = _find_target_table(slide, "buyers")
    if table is None:
        try:
            print("[PPT][TRACE] Buyers table NOT found", flush=True)
        except Exception:
            pass
        return False
    if len(table.columns) < 8:
        # Expect 8 columns: Rank | Name | Category | Shares Acquired | (2) current | (2) previous
        return False

    # Resolve column mappings from DF
    rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Name of Holder", "Shareholder Name", "Institution"]) or None
    cat_label_col = _ci_col(df, ["Category Label", "Category"]) or None
    buy_shares_col = _ci_col(df, ["Buy Shares", "Shares Acquired during the Week", "Shares Bought", "Shares Acquired"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])

    # Determine header rows count
    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(7, len(table.columns)))) else 1

    # Infer which side is current/previous by parsing header row dates (expected at col 4 and 6)
    left_is_prev = False if _force_current_left() else True
    if not _force_current_left():
        try:
            hdr_left_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
            hdr_right_txt = (table.cell(0, 6).text or "").strip() if len(table.columns) > 6 else ""
            dl = _parse_date_label_to_dt(hdr_left_txt)
            dr = _parse_date_label_to_dt(hdr_right_txt)
            if dl and dr:
                left_is_prev = dl < dr
        except Exception:
            pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 4, 6, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 4, 6, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    # Prepare to write
    start_row = header_rows
    # Ensure table has exactly header_rows + len(df) rows so extra blank rows don't stretch height
    _ensure_table_data_row_capacity(table, len(df), header_rows=header_rows)
    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0

    # Clear all data rows first
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    for i in range(n):
        row = df.iloc[i]
        r = start_row + i
        # Col 0: Rank (center)
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        # Col 1: Shareholder Name (left, shrink) — enforce 10pt baseline like Institutional
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)
        # Col 2: Category (left)
        _write_data_cell(table, r, 2, (row[cat_label_col] if cat_label_col in df.columns else ""), align=PP_ALIGN.LEFT, font_size_pt=10, word_wrap=False)
        # Col 3: Shares Acquired during the Week (center)
        _write_data_cell(table, r, 3, (row[buy_shares_col] if buy_shares_col in df.columns else ""), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Determine date block positions
        if left_is_prev:
            prev_block_cols = (4, 5)
            curr_block_cols = (6, 7)
        else:
            curr_block_cols = (4, 5)
            prev_block_cols = (6, 7)

        # Prev block values
        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Curr block values
        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    # Enforce uniform 10pt font on entire table after writing data.
    # Also compact spacing and constrain row/shape heights so all rows fit on the slide.
    try:
        print(f"[PPT][TRACE] Buyers applying heights: rows={len(table.rows)}, cols={len(table.columns)}, header_rows={header_rows}", flush=True)
    except Exception:
        pass
    # Match the manually-fixed Buyers slide: fixed table box + no text-to-fit.
    _disable_table_autofit(table, start_row=0)
    _compact_data_rows(table, 0)
    footer_top = _find_footer_legend_top(slide)
    if footer_top is not None:
        top_emu = Inches(1.145)
        gap_emu = Inches(0.02)
        target_h_emu = max(0, int(footer_top) - int(top_emu) - int(gap_emu))
        _set_table_shape_top(slide, table, top_inches=1.145)
        _set_table_shape_height(slide, table, height_inches=float(target_h_emu) / 914400.0)
        _fit_table_rows_to_shape_height(table, target_h_emu, header_rows=header_rows)
    else:
        _set_table_shape_top(slide, table, top_inches=1.145)
        _set_table_shape_height(slide, table, height_inches=5.95)
    _enforce_table_font_size(table, 0, font_size_pt=10)
    _enforce_column_font_size(table, 1, start_row=header_rows, font_size_pt=9)

    return True

def update_sellers_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 20 Sellers' table on the template slide with DB data (update-only)."""
    try:
        print("[PPT][TRACE] Enter update_sellers_table_on_template", flush=True)
    except Exception:
        pass
    # Locate slide
    slide = _find_slide_by_title(prs, "Top 20 Sellers")
    if slide is None:
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            if "sellers" in (" ".join(shp.text_frame.text.split()).strip().lower()):
                                slide = s
                                break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = (shp.text_frame.text or "").strip().lower()
                        if "top 20 sellers" in txt or "sellers" in txt:
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None and len(prs.slides) > 3:
        slide = prs.slides[3]
    if slide is None:
        return False

    table = _find_target_table(slide, "sellers")
    if table is None:
        try:
            print("[PPT][TRACE] Sellers table NOT found", flush=True)
        except Exception:
            pass
        return False
    if len(table.columns) < 8:
        return False

    # Column mappings
    rank_col = _ci_col(df, ["serial.no", "Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Name of Holder", "Shareholder Name", "Institution"]) or None
    cat_label_col = _ci_col(df, ["Category Label", "Category"]) or None
    sold_shares_col = _ci_col(df, ["sold Shares", "Shares Sold during the Week", "Shares Sold"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])

    # Header rows
    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(7, len(table.columns)))) else 1

    # Determine date block side by header dates at cols 4 and 6
    left_is_prev = False if _force_current_left() else True
    if not _force_current_left():
        try:
            hdr_left_txt = (table.cell(0, 4).text or "").strip() if len(table.columns) > 4 else ""
            hdr_right_txt = (table.cell(0, 6).text or "").strip() if len(table.columns) > 6 else ""
            dl = _parse_date_label_to_dt(hdr_left_txt)
            dr = _parse_date_label_to_dt(hdr_right_txt)
            if dl and dr:
                left_is_prev = dl < dr
        except Exception:
            pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 4, 6, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 4, 6, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    start_row = header_rows
    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0

    # Clear data rows
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    for i in range(n):
        row = df.iloc[i]
        r = start_row + i
        # Rank
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        # Name — enforce 10pt baseline like Institutional
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)
        # Category
        _write_data_cell(table, r, 2, (row[cat_label_col] if cat_label_col in df.columns else ""), align=PP_ALIGN.LEFT, font_size_pt=10, word_wrap=False)
        # Shares Sold during the Week
        _write_data_cell(table, r, 3, (row[sold_shares_col] if sold_shares_col in df.columns else ""), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        if left_is_prev:
            prev_block_cols = (4, 5)
            curr_block_cols = (6, 7)
        else:
            curr_block_cols = (4, 5)
            prev_block_cols = (6, 7)

        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    # Enforce uniform 10pt font on entire table after writing data.
    # Also compact spacing and constrain row/shape heights so all rows fit on the slide.
    try:
        print(f"[PPT][TRACE] Sellers applying heights: rows={len(table.rows)}, cols={len(table.columns)}, header_rows={header_rows}", flush=True)
    except Exception:
        pass
    # Match the manually-fixed Buyers slide: fixed table box + no text-to-fit.
    _disable_table_autofit(table, start_row=0)
    _compact_data_rows(table, 0)
    footer_top = _find_footer_legend_top(slide)
    if footer_top is not None:
        top_emu = Inches(1.145)
        gap_emu = Inches(0.02)
        target_h_emu = max(0, int(footer_top) - int(top_emu) - int(gap_emu))
        _set_table_shape_top(slide, table, top_inches=1.145)
        _set_table_shape_height(slide, table, height_inches=float(target_h_emu) / 914400.0)
        _fit_table_rows_to_shape_height(table, target_h_emu, header_rows=header_rows)
    else:
        _set_table_shape_top(slide, table, top_inches=1.145)
        _set_table_shape_height(slide, table, height_inches=5.95)
    _enforce_table_font_size(table, 0, font_size_pt=10)
    _enforce_column_font_size(table, 1, start_row=header_rows, font_size_pt=9)
    return True

def _normalize_shareholder_name(name: str) -> str:
    """Normalize shareholder name casing conservatively.
    - Title-case general words
    - Preserve common acronyms and short uppercase abbreviations
    - Preserve camel-cased tokens (e.g., iShares)
    - Handle hyphen and slash separated sub-tokens
    """
    if name is None:
        return ""
    s = str(name).strip()
    if not s:
        return s

    keep_upper = {
        "FII", "FPI", "DII", "MF", "MFs", "ETF", "ETFs", "PF", "AIF",
        "LLP", "PTE", "LTD", "PLC", "INC", "LLC", "PMS", "SWF", "IF",
        "SBI", "BNP", "ICICI", "UTI", "HDFC", "LIC", "QOG", "GQG", "HSBC",
    }

    def _tc_token(tok: str) -> str:
        if not tok or tok.isdigit():
            return tok
        core = tok.strip()
        # Preserve tokens already mixed-case (camelCase) or containing digits
        if any(ch.islower() for ch in core) and any(ch.isupper() for ch in core):
            return core
        # Preserve known acronyms/casings
        if core in keep_upper:
            return core
        if core.upper() in keep_upper:
            return core.upper()
        if core.isalpha() and core.isupper() and len(core) <= 3:
            return core
        return core.capitalize()

    def _process_separators(token: str, seps=("-", "/")) -> str:
        # Title-case subparts split by separators, keep separators
        out = token
        for sep in seps:
            if sep in out:
                parts = [ _tc_token(p) for p in out.split(sep) ]
                out = sep.join(parts)
        return _tc_token(out) if all(sep not in token for sep in seps) else out

    # Process word by word, preserving leading/trailing punctuation
    tokens = s.split(" ")
    out_tokens = []
    for i, tok in enumerate(tokens):
        # Strip surrounding punctuation but keep it in output
        prefix = ""
        suffix = ""
        lead_chars = "([{\"'"
        trail_chars = ')]}",.:;'
        while tok and tok[0] in lead_chars:
            prefix += tok[0]
            tok = tok[1:]
        while tok and tok[-1] in trail_chars:
            suffix = tok[-1] + suffix
            tok = tok[:-1]
        new_core = _process_separators(tok)
        # Lowercase small connector words except first token
        small = {"of", "and", "the", "in", "for", "on", "at", "by", "to"}
        if i > 0 and new_core.lower() in small:
            new_core = new_core.lower()
        out_tokens.append(prefix + new_core + suffix)
    return " ".join(out_tokens)

def _write_data_cell(table, r: int, c: int, value, align=PP_ALIGN.CENTER, font_size_pt: Optional[int] = 10, allow_shrink: bool = False, word_wrap: Optional[bool] = True, font_color: Optional[RGBColor] = None):
    """Write data text into a table cell with Arial font and wrapping.
    - font_size_pt: base font size in points (default 10). If None and allow_shrink=True, size is chosen based on text length.
    - allow_shrink: when True, reduce the font slightly for long strings to avoid overflow.
    """
    try:
        if not (0 <= r < len(table.rows) and 0 <= c < len(table.columns)):
            return False
        cell = table.cell(r, c)
        # Preserve the existing paragraph (and its spacing) and only replace runs
        tf = cell.text_frame
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        # Respect template's wrapping unless explicitly specified
        if word_wrap is not None:
            tf.word_wrap = bool(word_wrap)
        # Ensure a single paragraph; extra/empty paragraphs can force minimum row height
        try:
            paras = list(getattr(tf, "paragraphs", []))
            if len(paras) > 1:
                for p in paras[1:]:
                    try:
                        p._p.getparent().remove(p._p)
                    except Exception:
                        pass
        except Exception:
            pass
        para = tf.paragraphs[0]
        # Clear existing runs but keep paragraph node so spacing/margins persist
        try:
            for old_run in list(getattr(para, "runs", [])):
                try:
                    old_run._r.getparent().remove(old_run._r)
                except Exception:
                    pass
        except Exception:
            pass
        para.alignment = align
        # Do not alter vertical anchor or paragraph spacing; keep template defaults
        run = para.add_run()
        try:
            # Collapse whitespace (including newlines) to keep rows single-line where possible
            run.text = "" if value is None else " ".join(str(value).split())
        except Exception:
            run.text = "" if value is None else str(value)
        run.font.name = "Arial"
        try:
            if r == 0:
                run.font.bold = True
        except Exception:
            pass
        if font_color is not None:
            try:
                run.font.color.rgb = font_color
            except Exception:
                pass
        # dynamic font sizing
        try:
            if allow_shrink and isinstance(run.text, str):
                ln = len(run.text)
                base = 10 if font_size_pt is None else font_size_pt
                size = base
                if ln > 65:
                    size = max(8, base - 2)
                elif ln > 45:
                    size = max(9, base - 1)
                else:
                    size = max(10 if base >= 10 else base, size)
                run.font.size = Pt(size)
            else:
                if font_size_pt is not None:
                    run.font.size = Pt(font_size_pt)
        except Exception:
            pass
        return True
    except Exception:
        return False

def _ci_col(df: pd.DataFrame, names):
    """Case-insensitive exact match for any of the candidate names in df.columns."""
    if df is None or df.empty:
        return None
    if isinstance(names, str):
        names = [names]
    lower_map = {str(c).strip().lower(): c for c in df.columns}
    for n in names:
        key = str(n).strip().lower()
        if key in lower_map:
            return lower_map[key]
    return None

def _find_two_latest_date_columns(df: pd.DataFrame):
    """Find two most recent date-named columns like mm/dd/YYYY in df; returns (prev_col, curr_col)."""
    # New schema support: stable Current/Previous columns in append-only tables.
    try:
        cols_lower = {str(c).strip().lower(): c for c in df.columns}
        if "current" in cols_lower and "previous" in cols_lower:
            return cols_lower["previous"], cols_lower["current"]
    except Exception:
        pass

    date_cols = []
    for c in df.columns:
        s = str(c).strip()
        # Fast-path: Accept mm/dd/YYYY or mm/dd/YY
        if re.match(r"^\d{1,2}/\d{1,2}/\d{2,4}$", s):
            try:
                from datetime import datetime as _dt
                # Try common US format first, then D/M
                for fmt in ("%m/%d/%Y", "%m/%d/%y", "%d/%m/%Y", "%d/%m/%y"):
                    try:
                        dt = _dt.strptime(s, fmt)
                        date_cols.append((dt, c))
                        break
                    except Exception:
                        continue
            except Exception:
                pass
            continue

        # General-path: try parsing column name like '06-Mar-26' / '06 Mar 2026' etc.
        try:
            dt = _parse_date_label_to_dt(s)
            if dt is not None:
                date_cols.append((dt, c))
        except Exception:
            pass
    if len(date_cols) < 2:
        return None, None
    date_cols.sort(key=lambda x: x[0])
    prev_col = date_cols[-2][1]
    curr_col = date_cols[-1][1]
    return prev_col, curr_col

def _parse_date_label_to_dt(s: str):
    """Best-effort parse for header date labels like '26-Dec-25' or '12/26/2025'. Returns datetime or None."""
    try:
        s = str(s or "").strip()
        if not s:
            return None
        for fmt in (
            "%d-%b-%y", "%d-%b-%Y",
            "%d %b %y", "%d %b %Y",
            "%m/%d/%Y", "%m/%d/%y",
            "%d/%m/%Y", "%d/%m/%y",
        ):
            try:
                return datetime.strptime(s, fmt)
            except Exception:
                continue
        return None
    except Exception:
        return None


def _format_header_date_label(value) -> str:
    """Format a dataframe/header date value into PPT label style like '26-Dec-25'."""
    try:
        s0 = str(value or "").strip()
        if s0.lower() in ("current", "curr"):
            return _format_header_date_label(CURRENT_WEEK_DATE)
        if s0.lower() in ("previous", "prev"):
            return _format_header_date_label(PREVIOUS_WEEK_DATE)
    except Exception:
        pass
    try:
        dt = _parse_date_label_to_dt(value)
        if dt is not None:
            return dt.strftime("%d-%b-%y")
    except Exception:
        pass
    s = str(value or "").strip()
    return s


def _apply_header_date_labels(table, left_col: int, right_col: int, left_value, right_value):
    """Write visible date labels into the first header row of an existing template table."""
    try:
        if len(table.rows) < 1:
            return False
        if 0 <= left_col < len(table.columns):
            _write_data_cell(table, 0, left_col, _format_header_date_label(left_value), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=HEADER_FONT_COLOR)
        if 0 <= right_col < len(table.columns):
            _write_data_cell(table, 0, right_col, _format_header_date_label(right_value), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=HEADER_FONT_COLOR)
        return True
    except Exception:
        return False


def _force_current_left() -> bool:
    raw = (_get_env("WSHP_FORCE_CURRENT_LEFT", "1") or "1").strip().lower()
    return raw not in ("0", "false", "no", "off")


def _sort_df_rankwise(df: pd.DataFrame) -> pd.DataFrame:
    try:
        if df is None or df.empty:
            return df
    except Exception:
        return df
    try:
        rank_col = _ci_col(df, [
            "rank",
            "serial.no",
            "serial no",
            "serial. no",
            "sr.no",
            "sr no",
            "sr. no",
            "s.no",
            "sno",
        ])
        if not rank_col:
            return df
        out = df.copy()
        out["__rank_sort__"] = pd.to_numeric(out[rank_col], errors="coerce")
        out = out.sort_values(by=["__rank_sort__"], ascending=True, na_position="last")
        out = out.drop(columns=["__rank_sort__"], errors="ignore")
        return out
    except Exception:
        return df

def update_institutional_table_on_template(prs: Presentation, df: pd.DataFrame):
    """Update the 'Top 20 Institutional Shareholders' table on the template slide with DB data."""
    slide = _find_slide_by_title(prs, "Top 20 Institutional Shareholders")
    if slide is None:
        # Fallback A: any slide whose TITLE placeholder contains 'institutional'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if getattr(shp, "is_placeholder", False):
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.TITLE and hasattr(shp, "text_frame"):
                            if "institutional" in (" ".join(shp.text_frame.text.split()).strip().lower()):
                                slide = s
                                break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        # Fallback B: scan any text shape on slides for 'Top 20 Institutional'
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = (shp.text_frame.text or "").strip().lower()
                        if "top 20 institutional" in txt or "institutional shareholders" in txt:
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None and len(prs.slides) > 1:
        # Fallback C: default to the 2nd slide per template instruction
        slide = prs.slides[1]
    if slide is None:
        return False
    table = _get_best_table(slide) or _get_first_table(slide)
    if table is None:
        return False
    # Expect at least 8 columns for this table
    if len(table.columns) < 7:
        return False

    # Resolve column mappings from DF
    rank_col = _ci_col(df, ["Serial. No", "Serial No", "Sr.No", "Sr. No", "Sr No", "Rank"]) or None
    name_col = _ci_col(df, ["Institution", "Shareholder Name"]) or None
    cat_label_col = _ci_col(df, ["Category Label"]) or None
    prev_hold_col, curr_hold_col = _find_two_latest_date_columns(df)
    prev_pct_col = _ci_col(df, ["% of Sh. Cap (previous)", "% of Sh. Cap (Previous)"])
    curr_pct_col = _ci_col(df, ["% of Sh. Cap (current)", "% of Sh. Cap (Current)"])
    change_col = _ci_col(df, ["MoM change in holdings", "MoM change in holding", "Change in Holding", "Change"])
    buy_pps_col  = _ci_col(df, ["Buy"]) or None
    sell_pps_col = _ci_col(df, ["Sell"]) or None

    # Determine header rows count (1 or 2). If second row has any text in header area, assume 2.
    header_rows = 2 if any(len((table.cell(1, c).text or "").strip()) > 0 for c in range(min(6, len(table.columns)))) else 1
    # Default mapping: assume current date block is on the LEFT like the template screenshot
    # We'll refine using header date labels when available
    left_is_prev = not _force_current_left()
    # Infer from top-level headers if available by parsing dates
    try:
        hdr_left_txt = (table.cell(0, 3).text or "").strip() if len(table.columns) > 3 else ""
        hdr_right_txt = (table.cell(0, 5).text or "").strip() if len(table.columns) > 5 else ""
        dl = _parse_date_label_to_dt(hdr_left_txt)
        dr = _parse_date_label_to_dt(hdr_right_txt)
        if dl and dr:
            # If the left header date is earlier than the right header date, then left is PREVIOUS
            left_is_prev = dl < dr
    except Exception:
        pass
    try:
        if left_is_prev:
            _apply_header_date_labels(table, 3, 5, prev_hold_col, curr_hold_col)
        else:
            _apply_header_date_labels(table, 3, 5, curr_hold_col, prev_hold_col)
    except Exception:
        pass

    # Write rows (preserve formatting). Assume 2 header rows.
    start_row = header_rows
    max_rows = len(table.rows) - start_row
    try:
        df = _sort_df_rankwise(df)
    except Exception:
        pass
    n = min(max_rows, len(df)) if max_rows > 0 else 0
    # Clear all data rows first
    for r in range(start_row, len(table.rows)):
        for c in range(len(table.columns)):
            try:
                table.cell(r, c).text = ""
            except Exception:
                pass

    for i in range(n):
        row = df.iloc[i]
        r = start_row + i
        # Col 0: Rank
        val_rank = row[rank_col] if (rank_col is not None and rank_col in df.columns) else (i + 1)
        _write_data_cell(table, r, 0, val_rank, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        # Col 1: Shareholder Name
        name_val = row[name_col] if name_col in df.columns else ""
        name_val = _normalize_shareholder_name(name_val)
        _write_data_cell(table, r, 1, name_val, align=PP_ALIGN.LEFT, font_size_pt=9, allow_shrink=False, word_wrap=False)
        # Col 2: Category (Category Label)
        _write_data_cell(table, r, 2, (row[cat_label_col] if cat_label_col in df.columns else ""), align=PP_ALIGN.LEFT, font_size_pt=10, word_wrap=False)

        # Determine positions
        if left_is_prev:
            prev_block_cols = (3, 4)
            curr_block_cols = (5, 6)
        else:
            # Left is current
            curr_block_cols = (3, 4)
            prev_block_cols = (5, 6)

        # Prev block values
        prev_hold_val = _format_holding(row[prev_hold_col]) if prev_hold_col in df.columns else ""
        prev_pct_val = _format_pct(row[prev_pct_col]) if prev_pct_col in df.columns else ""
        _write_data_cell(table, r, prev_block_cols[0], prev_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, prev_block_cols[1], prev_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Curr block values
        curr_hold_val = _format_holding(row[curr_hold_col]) if curr_hold_col in df.columns else ""
        curr_pct_val = _format_pct(row[curr_pct_col]) if curr_pct_col in df.columns else ""
        _write_data_cell(table, r, curr_block_cols[0], curr_hold_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
        _write_data_cell(table, r, curr_block_cols[1], curr_pct_val, align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

        # Change in Holding (pps): prefer Buy (green), else Sell (red), else fallback to change_col
        if len(table.columns) > 7:
            def _nonzero(x):
                try:
                    s = str(x).strip()
                    if not s or s.lower() == 'none':
                        return False
                    v = float(s.replace(',', '').replace('(', '-').replace(')', ''))
                    return abs(v) > 0
                except Exception:
                    return bool(str(x).strip())
            buy_v = (row[buy_pps_col] if (buy_pps_col in df.columns) else None) if buy_pps_col is not None else None
            sell_v = (row[sell_pps_col] if (sell_pps_col in df.columns) else None) if sell_pps_col is not None else None
            if _nonzero(buy_v):
                _write_data_cell(table, r, 7, _format_change_value(buy_v), align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)
            elif _nonzero(sell_v):
                _write_data_cell(table, r, 7, f"({_format_change_value(sell_v)})", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False, font_color=RED_COLOR)
            else:
                _write_data_cell(table, r, 7, "-", align=PP_ALIGN.CENTER, font_size_pt=10, word_wrap=False)

    return True

def build_title_slide(prs, bu_name: Optional[str] = None):
    """Update only the existing standalone date text on slide 1, preserving title layout."""
    slide = prs.slides[0] if len(prs.slides) > 0 else add_blank_slide(prs)

    def _looks_like_date_text(text: str) -> bool:
        try:
            s = " ".join(str(text or "").split()).strip()
            if not s:
                return False
            return _parse_date_label_to_dt(s) is not None
        except Exception:
            return False

    try:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    tf = shape.text_frame
                    existing = tf.text or ""
                    if _looks_like_date_text(existing):
                        tf.clear()
                        p = tf.paragraphs[0]
                        r = p.add_run()
                        r.text = str(CURRENT_WEEK_DATE or "")
                        try:
                            r.font.name = "Arial"
                        except Exception:
                            pass
                        break
            except Exception:
                continue
    except Exception:
        pass

    try:
        if bu_name:
            for shape in slide.shapes:
                try:
                    if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                        continue
                    tf = shape.text_frame
                    txt = (tf.text or "").strip()
                    low = " ".join(txt.split()).lower()
                    if "weekly" in low and "shareholder" in low and "movement" in low:
                        new_txt = f"{bu_name} | Weekly Shareholder Movement"
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        try:
                            # Keep the title on one line; allow auto-fit to shrink text if needed.
                            tf.word_wrap = False
                        except Exception:
                            pass
                        try:
                            # Some templates ignore TEXT_TO_FIT_SHAPE for placeholders.
                            # Use explicit font sizing and keep auto-size disabled.
                            tf.auto_size = MSO_AUTO_SIZE.NONE
                        except Exception:
                            pass
                        try:
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                        except Exception:
                            pass
                        try:
                            p = tf.paragraphs[0]
                        except Exception:
                            continue
                        try:
                            p.alignment = PP_ALIGN.RIGHT
                        except Exception:
                            pass
                        r = p.add_run()
                        r.text = new_txt
                        try:
                            r.font.name = "Arial"
                        except Exception:
                            pass
                        try:
                            # Heuristic sizing: shrink if the string is long enough to wrap.
                            ln = len(new_txt)
                            sz = 18
                            if ln >= 85:
                                sz = 11
                            elif ln >= 70:
                                sz = 12
                            elif ln >= 55:
                                sz = 13
                            elif ln >= 45:
                                sz = 14
                            r.font.size = Pt(sz)
                        except Exception:
                            pass

                        # Enforce the font settings across all runs to avoid theme/master overrides.
                        try:
                            for pp in tf.paragraphs:
                                try:
                                    pp.alignment = PP_ALIGN.RIGHT
                                except Exception:
                                    pass
                                for rr in pp.runs:
                                    try:
                                        rr.font.name = "Arial"
                                    except Exception:
                                        pass
                                    try:
                                        rr.font.size = Pt(sz)
                                    except Exception:
                                        pass
                        except Exception:
                            pass
                        break
                except Exception:
                    continue
    except Exception:
        pass

    try:
        def _norm(s: str) -> str:
            return " ".join(str(s or "").split()).strip().lower()

        for shp in list(slide.shapes):
            try:
                # PowerPoint may render slide-number placeholders at open time even if
                # text_frame.text is empty here; remove by placeholder type as well.
                if getattr(shp, "is_placeholder", False):
                    try:
                        phf = shp.placeholder_format
                        if hasattr(phf, "type") and phf.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                            try:
                                slide.shapes._spTree.remove(shp._element)
                                continue
                            except Exception:
                                pass
                    except Exception:
                        pass

                try:
                    nm = _norm(getattr(shp, "name", ""))
                    if "slide number" in nm or nm == "slidenumber":
                        try:
                            slide.shapes._spTree.remove(shp._element)
                            continue
                        except Exception:
                            pass
                except Exception:
                    pass

                if getattr(shp, "has_text_frame", False):
                    txt = _norm(shp.text_frame.text)
                    if txt in ("slidenumber", "slide number"):
                        try:
                            slide.shapes._spTree.remove(shp._element)
                        except Exception:
                            pass
            except Exception:
                continue
    except Exception:
        pass

    return slide


def update_toc_on_template(prs, toc_items):
    """Find the existing 'Table of Contents' slide in the template and overwrite
    its list with the provided toc_items. Tries to update an existing table; if
    none exists, adds a new 2-column table to that slide.
    """
    # Locate the TOC slide by title or by scanning text
    slide = _find_slide_by_title(prs, "Table of Contents")
    if slide is None:
        for s in prs.slides:
            try:
                for shp in s.shapes:
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = (shp.text_frame.text or "").strip().lower()
                        if "table of contents" in " ".join(txt.split()):
                            slide = s
                            break
                if slide is not None:
                    break
            except Exception:
                continue
    if slide is None:
        # Fallback: assume slide 2 is TOC if present
        try:
            if len(prs.slides) > 1:
                slide = prs.slides[1]
        except Exception:
            slide = None
    if slide is None:
        return False

    def _norm_txt(s: str) -> str:
        return " ".join(str(s or "").split()).strip().lower()

    def _set_tf_text_preserve_format(tf, new_text: str) -> None:
        try:
            if not tf.paragraphs:
                tf.text = new_text
                return
            p = tf.paragraphs[0]
            if not p.runs:
                tf.text = new_text
                return
            p.runs[0].text = new_text
            for r in p.runs[1:]:
                r.text = ""
        except Exception:
            try:
                tf.text = new_text
            except Exception:
                pass

    try:
        # 1) Prefer updating an existing table (keeps template formatting)
        for shp in slide.shapes:
            try:
                if getattr(shp, "has_table", False) and shp.table is not None:
                    tbl = shp.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    if cols < 2:
                        continue
                    max_rows = min(rows, len(toc_items))
                    for i in range(max_rows):
                        _set_tf_text_preserve_format(tbl.cell(i, 0).text_frame, str(i + 1))
                        _set_tf_text_preserve_format(tbl.cell(i, 1).text_frame, str(toc_items[i]))
                    for i in range(max_rows, rows):
                        _set_tf_text_preserve_format(tbl.cell(i, 0).text_frame, "")
                        _set_tf_text_preserve_format(tbl.cell(i, 1).text_frame, "")
                    return True
            except Exception:
                continue

        # 2) Fallback: update a textbox list (most paragraphs) without changing layout
        candidate = None
        best_score = -1
        for shp in slide.shapes:
            try:
                if not getattr(shp, "has_text_frame", False):
                    continue
                txt = _norm_txt(shp.text_frame.text)
                if not txt or "table of contents" in txt:
                    continue
                score = len(shp.text_frame.paragraphs)
                if score > best_score:
                    best_score = score
                    candidate = shp
            except Exception:
                continue

        if candidate is None:
            return False

        tf = candidate.text_frame
        for i, item in enumerate(toc_items):
            if i >= len(tf.paragraphs):
                break
            p = tf.paragraphs[i]
            if not p.runs:
                try:
                    p.text = f"{i + 1} {item}"
                except Exception:
                    continue
            else:
                p.runs[0].text = f"{i + 1} {item}"
                for r in p.runs[1:]:
                    r.text = ""
        for j in range(len(toc_items), len(tf.paragraphs)):
            try:
                p = tf.paragraphs[j]
                if p.runs:
                    p.runs[0].text = ""
                    for r in p.runs[1:]:
                        r.text = ""
                else:
                    p.text = ""
            except Exception:
                continue
        return True
    except Exception:
        return False

def build_toc_slide(prs):
    """Build the Table of Contents slide."""
    slide = add_blank_slide(prs)

    # Title: "Table of Contents"
    txBox = slide.shapes.add_textbox(
        Emu(5657849), Emu(1423331), Emu(3620905), Emu(747837)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Table of Contents"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # TOC items
    toc_items = [
        "Top 20 Institutional Shareholders",
        "Top 20 Buyers",
        "Top 20 Sellers",
        "New Entry / Exits",
        "Top 10 FII’s & FPI’s",
        "Top 10 MF’s",
        "Top 10 Insurance & PF’s",
        "Top 10 AIF’s",
    ]

    num_rows = len(toc_items)
    num_cols = 2
    table_left = Emu(5779770)
    table_top = Emu(2284098)
    table_width = Emu(6006327)
    row_height = Emu(335931)
    table_height = row_height * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top,
                                         table_width, table_height)
    table = table_shape.table
    remove_table_styling(table)

    # Column widths matching template
    table.columns[0].width = Emu(622658)
    table.columns[1].width = Emu(5383669)

    for i, item in enumerate(toc_items):
        # Number column
        set_cell_style(
            table.cell(i, 0), str(i + 1),
            font_size=Pt(14), bold=False,
            font_color=RGBColor(0x1B, 0x3A, 0x5C),
            alignment=PP_ALIGN.CENTER, font_name="Arial"
        )
        # Item name column
        set_cell_style(
            table.cell(i, 1), item,
            font_size=Pt(14), bold=False,
            font_color=RGBColor(0x1B, 0x3A, 0x5C),
            alignment=PP_ALIGN.LEFT, font_name="Arial"
        )
        # Add borders
        for c in range(num_cols):
            set_cell_border(table.cell(i, c))

    return slide


def _build_7col_header(table, date_current, date_previous, change_label="Change in Holding"):
    """Build a 2-row merged header for 7-column tables:
    Row 0: Rank | Shareholder Name | <current_date> (merged 2) | <prev_date> (merged 2) | Change in Holding
    Row 1: (merged) | (merged)     | Holding | % of Share Capital | Holding | % of Share Capital | (merged)
    """
    def _fmt_date_label(s):
        s = str(s or "").strip()
        for fmt in ("%d-%b-%y", "%d-%b-%Y", "%d/%m/%Y", "%m/%d/%Y", "%d-%m-%Y", "%d.%m.%Y"):
            try:
                dt = datetime.strptime(s, fmt)
                return dt.strftime("%d-%b-%y")
            except Exception:
                continue
        return s

    date_current = _fmt_date_label(date_current)
    date_previous = _fmt_date_label(date_previous)
    # Row 0
    set_cell_style(table.cell(0, 0), "Rank", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 1), "Shareholder Name", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR, alignment=PP_ALIGN.LEFT)
    set_cell_style(table.cell(0, 2), date_current, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 3), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 4), date_previous, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 5), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 6), change_label, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    # Row 1
    set_cell_style(table.cell(1, 0), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 1), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 2), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 3), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 4), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 5), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 6), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    # Merge: Rank (rows 0-1), Shareholder Name (rows 0-1), Change (rows 0-1)
    merge_cells(table, 0, 0, 1, 0)  # Rank
    merge_cells(table, 0, 1, 1, 1)  # Shareholder Name
    merge_cells(table, 0, 2, 0, 3)  # Current date cols
    merge_cells(table, 0, 4, 0, 5)  # Previous date cols
    merge_cells(table, 0, 6, 1, 6)  # Change in Holding


def _build_8col_header_institutional(table, date_current, date_previous):
    """Build 2-row merged header for 8-column institutional table:
    Row 0: Rank | Shareholder Name | Category | <current_date> (merged 2) | <prev_date> (merged 2) | Change in Holding
    """
    def _fmt_date_label(s):
        s = str(s or "").strip()
        for fmt in ("%d-%b-%y", "%d-%b-%Y", "%d/%m/%Y", "%m/%d/%Y", "%d-%m-%Y", "%d.%m.%Y"):
            try:
                dt = datetime.strptime(s, fmt)
                return dt.strftime("%d-%b-%y")
            except Exception:
                continue
        return s

    date_current = _fmt_date_label(date_current)
    date_previous = _fmt_date_label(date_previous)
    set_cell_style(table.cell(0, 0), "Rank", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 1), "Shareholder Name", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR, alignment=PP_ALIGN.LEFT)
    set_cell_style(table.cell(0, 2), "Category", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 3), date_current, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 4), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 5), date_previous, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 6), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 7), "Change in Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    set_cell_style(table.cell(1, 0), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 1), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 2), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 3), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 4), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 5), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 6), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 7), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    merge_cells(table, 0, 0, 1, 0)  # Rank
    merge_cells(table, 0, 1, 1, 1)  # Shareholder Name
    merge_cells(table, 0, 2, 1, 2)  # Category
    merge_cells(table, 0, 3, 0, 4)  # Current date cols
    merge_cells(table, 0, 5, 0, 6)  # Previous date cols
    merge_cells(table, 0, 7, 1, 7)  # Change in Holding


def _build_8col_header_buyers_sellers(table, date_current, date_previous, action_col_name):
    """Build 2-row merged header for 8-column buyers/sellers table:
    Row 0: Rank | Shareholder Name | Category | <action> | <current_date> (merged 2) | <prev_date> (merged 2)
    """
    def _fmt_date_label(s):
        s = str(s or "").strip()
        for fmt in ("%d-%b-%y", "%d-%b-%Y", "%d/%m/%Y", "%m/%d/%Y", "%d-%m-%Y", "%d.%m.%Y"):
            try:
                dt = datetime.strptime(s, fmt)
                return dt.strftime("%d-%b-%y")
            except Exception:
                continue
        return s

    date_current = _fmt_date_label(date_current)
    date_previous = _fmt_date_label(date_previous)
    set_cell_style(table.cell(0, 0), "Rank", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 1), "Shareholder Name", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR, alignment=PP_ALIGN.LEFT)
    set_cell_style(table.cell(0, 2), "Category", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 3), action_col_name, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 4), date_current, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 5), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 6), date_previous, HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(0, 7), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    set_cell_style(table.cell(1, 0), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 1), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 2), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 3), "", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 4), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 5), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 6), "Holding", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)
    set_cell_style(table.cell(1, 7), "% of Share Capital", HEADER_FONT_SIZE, True, HEADER_FONT_COLOR, bg_color=HEADER_BG_COLOR)

    merge_cells(table, 0, 0, 1, 0)  # Rank
    merge_cells(table, 0, 1, 1, 1)  # Shareholder Name
    merge_cells(table, 0, 2, 1, 2)  # Category
    merge_cells(table, 0, 3, 1, 3)  # Action column
    merge_cells(table, 0, 4, 0, 5)  # Current date cols
    merge_cells(table, 0, 6, 0, 7)  # Previous date cols


def _fill_data_rows_8col(table, df, start_row=2, change_col_idx=7):
    """Fill data rows for 8-column tables.
    Expected df columns: rank, shareholder_name, category, col3, 
                          current_holding, current_pct, prev_holding, prev_pct
    The change_col_idx column gets colored green/red based on value.
    """
    for i, row in df.iterrows():
        r = start_row + i
        if r >= len(table.rows):
            break
        values = list(row)
        for c in range(min(len(values), 8)):
            val = str(values[c]) if values[c] is not None else ""
            if c == change_col_idx:
                color = get_change_color(val)
            else:
                color = DATA_FONT_COLOR
            align = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
            set_cell_style(table.cell(r, c), val, DATA_FONT_SIZE, False, color,
                          bg_color=WHITE_BG, alignment=align)
            set_cell_border(table.cell(r, c))


def _date_tokens(date_str: str):
    tokens = set()
    if not date_str:
        return tokens
    s = str(date_str).strip()
    if not s:
        return tokens
    tokens.add(s.lower())
    tokens.add(s.replace('/', '-').lower())
    for fmti in ("%d-%b-%y", "%d-%b-%Y", "%d/%m/%Y", "%m/%d/%Y", "%d-%m-%Y", "%d.%m.%Y"):
        try:
            dt = datetime.strptime(s, fmti)
            for fmto in ("%d-%b-%y", "%d-%b-%Y", "%d/%m/%Y", "%m/%d/%Y", "%d-%m-%Y"):
                tokens.add(dt.strftime(fmto).lower())
        except Exception:
            continue
    return tokens


def _pick_col(df, must_include_all):
    toks = [t.lower() for t in must_include_all if t]
    for col in df.columns:
        name = str(col).lower()
        if all(tok in name for tok in toks):
            return col
    return None


def _prepare_institutional_df(df, current_date_str, prev_date_str):
    if df is None or df.empty:
        return df
    # --- exact/ci-exact helpers ---
    def _exact_col(names):
        # Try exact then case-insensitive exact
        for n in names:
            if n in df.columns:
                return n
        lower_map = {str(c).strip().lower(): c for c in df.columns}
        for n in names:
            key = str(n).strip().lower()
            if key in lower_map:
                return lower_map[key]
        return None

    # Prefer exact names first as per user's specification
    rank_col = _exact_col(["Rank"]) or _pick_col(df, ["rank"]) or None
    inst_col = _exact_col(["Institution"]) or _pick_col(df, ["institution"]) or _pick_col(df, ["shareholder", "name"]) or None
    cat_col = _exact_col(["Category Label"]) or _pick_col(df, ["category", "label"]) or _pick_col(df, ["category"]) or None
    cur_pct_col = _exact_col(["% of Sh. Cap (Current)"]) or _pick_col(df, ["%", "sh.", "cap", "current"]) or _pick_col(df, ["percent", "current"]) or None
    prev_pct_col = _exact_col(["% of Sh. Cap (Previous)"]) or _pick_col(df, ["%", "sh.", "cap", "previous"]) or _pick_col(df, ["percent", "previous"]) or None
    cur_tokens = _date_tokens(current_date_str)
    prev_tokens = _date_tokens(prev_date_str)

    def _find_hold_for(tokens_set):
        # Try strong exact patterns first: 'Holdings - <date>'
        variants = []
        for t in tokens_set:
            # generate common display variants
            variants.extend([
                f"Holdings - {t}", f"Holding - {t}", f"{t} Holdings", f"{t} Holding"
            ])
        col = _exact_col(variants)
        if col:
            return col
        # Next, allow columns named directly by the date token (e.g., '26-Dec-25')
        # but avoid percentage columns.
        for c in df.columns:
            name = str(c).lower()
            if any(tok in name for tok in tokens_set):
                if not any(x in name for x in ["%", "percent", "sh.", "share", "cap", "capital"]):
                    return c
        # Fallback: substring match 'hold' + any date token
        for col in df.columns:
            name = str(col).lower()
            if "hold" in name and any(tok in name for tok in tokens_set):
                return col
        for hint in ("current", "curr"):
            c = _pick_col(df, ["hold", hint])
            if c:
                return c
        return None

    cur_hold_col = _find_hold_for(cur_tokens)
    prev_hold_col = _find_hold_for(prev_tokens)
    if not cur_hold_col or not prev_hold_col:
        holding_like = [c for c in df.columns if "hold" in str(c).lower()]
        if len(holding_like) >= 2:
            cur_hold_col = cur_hold_col or holding_like[0]
            prev_hold_col = prev_hold_col or holding_like[1]

    change_col = _pick_col(df, ["mom", "change"]) or _pick_col(df, ["change", "holding"]) or None

    cols_order = []
    if rank_col and rank_col in df.columns:
        cols_order.append(rank_col)
    else:
        cols_order.append(None)
    cols_order += [inst_col, cat_col, cur_hold_col, cur_pct_col, prev_hold_col, prev_pct_col, change_col]

    out_cols = []
    for c in cols_order:
        if c is None or c not in df.columns:
            out_cols.append(pd.Series([""] * len(df)))
        else:
            out_cols.append(df[c])

    out = pd.concat(out_cols, axis=1)
    out.columns = [
        "rank", "shareholder_name", "category", "current_holding", "current_pct",
        "prev_holding", "prev_pct", "change"
    ]
    if rank_col is None or rank_col not in df.columns:
        out["rank"] = range(1, len(out) + 1)
    return out

def _fill_data_rows_7col(table, df, start_row=2, change_col_idx=6):
    """Fill data rows for 7-column tables."""
    for i, row in df.iterrows():
        r = start_row + i
        if r >= len(table.rows):
            break
        values = list(row)
        for c in range(min(len(values), 7)):
            val = str(values[c]) if values[c] is not None else ""
            if c == change_col_idx:
                color = get_change_color(val)
            else:
                color = DATA_FONT_COLOR
            align = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
            set_cell_style(table.cell(r, c), val, DATA_FONT_SIZE, False, color,
                          bg_color=WHITE_BG, alignment=align)
            set_cell_border(table.cell(r, c))


def _add_title_textbox(slide, title_text):
    """Add a title textbox at the standard position."""
    txBox = slide.shapes.add_textbox(
        Emu(251200), Emu(402051), Emu(7729163), Emu(341632)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = title_text
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)


def _add_note_textbox(slide, note_text="Note: Shares in Lakhs"):
    """Add a 'Note: Shares in Lakhs' textbox at the standard position."""
    txBox = slide.shapes.add_textbox(
        Emu(10504688), Emu(820521), Emu(1381552), Emu(230832)
    )
    tf = txBox.text_frame
    para = tf.paragraphs[0]
    run1 = para.add_run()
    if ":" in note_text:
        parts = note_text.split(":", 1)
        run1.text = parts[0] + ": "
        run1.font.size = NOTE_FONT_SIZE
        run1.font.name = "Arial"
        run1.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
        run2 = para.add_run()
        run2.text = parts[1].strip()
        run2.font.size = NOTE_FONT_SIZE
        run2.font.name = "Arial"
        run2.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    else:
        run1.text = note_text
        run1.font.size = NOTE_FONT_SIZE
        run1.font.name = "Arial"
        run1.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)


def _add_footer_textbox(slide, footer_text):
    """Add abbreviation footer at the bottom."""
    txBox = slide.shapes.add_textbox(
        Emu(348016), Emu(6583163), Emu(11354766), Emu(215444)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]

    # Parse bold abbreviations: "FIIs: Foreign Institutional Investors; ..."
    parts = footer_text.split("; ")
    for idx, part in enumerate(parts):
        if ": " in part:
            abbr, full = part.split(": ", 1)
            run_bold = para.add_run()
            run_bold.text = abbr + ": "
            run_bold.font.size = FOOTER_FONT_SIZE
            run_bold.font.bold = True
            run_bold.font.name = "Arial"
            run_bold.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

            run_norm = para.add_run()
            suffix = "; " if idx < len(parts) - 1 else ""
            run_norm.text = full + suffix
            run_norm.font.size = FOOTER_FONT_SIZE
            run_norm.font.bold = False
            run_norm.font.name = "Arial"
            run_norm.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
        else:
            run = para.add_run()
            suffix = "; " if idx < len(parts) - 1 else ""
            run.text = part + suffix
            run.font.size = FOOTER_FONT_SIZE
            run.font.name = "Arial"
            run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)


# â”€â”€ Individual Slide Builders â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def build_top20_institutional_slide(prs, df):
    """Slide: Top 20 Institutional Shareholders (8 columns)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 20 Institutional Shareholders")
    _add_note_textbox(slide, "Shares in Lakhs")

    # Table: 22 rows x 8 cols (2 header + 20 data)
    num_data_rows = min(len(df), 20)
    total_rows = 2 + num_data_rows  # 2 header rows + data
    num_cols = 8

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1046123), Emu(11521440), Emu(5394960)
    )
    table = table_shape.table
    remove_table_styling(table)

    # Column widths
    col_widths = [457200, 4206240, 1155703, 1155703, 1155703, 1155703, 1155703, 1155703]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    # Header
    _build_8col_header_institutional(table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE)

    # Add borders to header
    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    # Data rows (map to expected structure/order)
    df_mapped = _prepare_institutional_df(df, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE)
    _fill_data_rows_8col(table, df_mapped, start_row=2, change_col_idx=7)

    # Footer
    _add_footer_textbox(slide,
        "FIIs: Foreign Institutional Investors; DIIs: Domestic Institutional Investors; "
        "MFs: Mutual Funds; IF: Index Fund; SWFs: Sovereign Wealth Funds; "
        "PMS: Portfolio Management Services")

    return slide


def build_top20_buyers_slide(prs, df):
    """Slide: Top 20 Buyers (8 columns with 'Shares Acquired' column)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 20 Buyers")
    _add_note_textbox(slide, "Note: Shares in Lakhs")

    num_data_rows = min(len(df), 20)
    total_rows = 2 + num_data_rows
    num_cols = 8

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1065524), Emu(11521440), Emu(5303520)
    )
    table = table_shape.table
    remove_table_styling(table)

    col_widths = [457200, 4206240, 1155703, 1155703, 1155703, 1155703, 1155703, 1155703]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    _build_8col_header_buyers_sellers(
        table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE,
        "Shares Acquired \nduring the Week"
    )

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    # Data rows (no dedicated "change" column for buyers â€” all black)
    _fill_data_rows_8col(table, df, start_row=2, change_col_idx=-1)

    _add_footer_textbox(slide,
        "FIIs: Foreign Institutional Investors; FPIs: Foreign Portfolio Investments; "
        "DIIs: Domestic Institutional Investors; MFs: Mutual Funds; "
        "IF: Index Fund; SWFs: Sovereign Wealth Funds; "
        "PMS: Portfolio Management Services")

    return slide


def build_top20_sellers_slide(prs, df):
    """Slide: Top 20 Sellers (8 columns with 'Shares Sold' column)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 20 Sellers")
    _add_note_textbox(slide, "Note: Shares in Lakhs")

    num_data_rows = min(len(df), 20)
    total_rows = 2 + num_data_rows
    num_cols = 8

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1065524), Emu(11521440), Emu(5303520)
    )
    table = table_shape.table
    remove_table_styling(table)

    col_widths = [457200, 4206240, 1155703, 1155703, 1155703, 1155703, 1155703, 1155703]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    _build_8col_header_buyers_sellers(
        table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE,
        "Shares Sold \nduring the Week"
    )

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_8col(table, df, start_row=2, change_col_idx=-1)

    _add_footer_textbox(slide,
        "FIIs: Foreign Institutional Investors; FPIs: Foreign Portfolio Investments; "
        "DIIs: Domestic Institutional Investors; MFs: Mutual Funds; "
        "IF: Index Fund; SWFs: Sovereign Wealth Funds; "
        "PMS: Portfolio Management Services")

    return slide


def build_top10_fii_fpi_slide(prs, df):
    """Slide: Top 10 FII's & FPI's (7 columns)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 20 FIIs & FPIs")
    _add_note_textbox(slide, "Shares in Lakhs")

    num_data_rows = min(len(df), 20)
    total_rows = 2 + num_data_rows
    num_cols = 7

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1046123), Emu(11521440), Emu(5486400)
    )
    table = table_shape.table
    remove_table_styling(table)

    col_widths = [548640, 4663440, 1280160, 1280160, 1280160, 1280160, 1280160]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    _build_7col_header(table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE)

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_7col(table, df, start_row=2, change_col_idx=6)

    _add_footer_textbox(slide,
        "FIIs: Foreign Institutional Investors; FPIs: Foreign Portfolio Investments")

    return slide


def build_top10_mf_slide(prs, df_active, df_passive):
    """Slide: Top 10 MFs with two sub-tables (Active + Passive)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 10 MFs")
    _add_note_textbox(slide, "Note: Shares in Lakhs")

    # â”€â”€ Active MF sub-header â”€â”€
    txBox_active = slide.shapes.add_textbox(
        Emu(348016), Emu(823344), Emu(11521440), Emu(258532)
    )
    tf_a = txBox_active.text_frame
    para_a = tf_a.paragraphs[0]
    run_a = para_a.add_run()
    run_a.text = "Top 10 Active MF Investors"
    run_a.font.size = Pt(12)
    run_a.font.bold = True
    run_a.font.name = "Arial"
    run_a.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Active MF Table (12 rows: 2 header + 10 data)
    num_active = min(len(df_active), 10)
    active_rows = 2 + num_active
    num_cols = 7

    table_shape_a = slide.shapes.add_table(
        active_rows, num_cols,
        Emu(348016), Emu(1082036), Emu(11521440), Emu(2560320)
    )
    table_a = table_shape_a.table
    remove_table_styling(table_a)

    col_widths = [548640, 4663440, 1280160, 1280160, 1280160, 1280160, 1280160]
    for i, w in enumerate(col_widths):
        table_a.columns[i].width = Emu(w)

    _build_7col_header(table_a, PREVIOUS_WEEK_DATE, CURRENT_WEEK_DATE)

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table_a.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_7col(table_a, df_active, start_row=2, change_col_idx=6)

    # â”€â”€ Passive MF sub-header â”€â”€
    txBox_passive = slide.shapes.add_textbox(
        Emu(348016), Emu(3670823), Emu(11521440), Emu(258532)
    )
    tf_p = txBox_passive.text_frame
    para_p = tf_p.paragraphs[0]
    run_p = para_p.add_run()
    run_p.text = "Top 10 Passive MF Investors"
    run_p.font.size = Pt(12)
    run_p.font.bold = True
    run_p.font.name = "Arial"
    run_p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Passive MF Table
    num_passive = min(len(df_passive), 10)
    passive_rows = 2 + num_passive

    table_shape_p = slide.shapes.add_table(
        passive_rows, num_cols,
        Emu(348016), Emu(3929512), Emu(11521440), Emu(2560320)
    )
    table_p = table_shape_p.table
    remove_table_styling(table_p)

    for i, w in enumerate(col_widths):
        table_p.columns[i].width = Emu(w)

    _build_7col_header(table_p, PREVIOUS_WEEK_DATE, CURRENT_WEEK_DATE)

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table_p.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_7col(table_p, df_passive, start_row=2, change_col_idx=6)

    # Footer
    _add_footer_textbox(slide, "MFs: Mutual Funds")

    return slide


def build_top10_insurance_pf_slide(prs, df):
    """Slide: Top 10 Insurance & PFs (7 columns)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 10 Insurance & PFs")
    _add_note_textbox(slide, "Shares in Lakhs")

    num_data_rows = min(len(df), 10)
    total_rows = 2 + num_data_rows
    num_cols = 7

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1046126), Emu(11521440), Emu(1682496)
    )
    table = table_shape.table
    remove_table_styling(table)

    col_widths = [548640, 4663440, 1280160, 1280160, 1280160, 1280160, 1280160]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    _build_7col_header(table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE,
                       "Change in Holding\n(pps)")

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_7col(table, df, start_row=2, change_col_idx=6)

    # Footer
    footer_txBox = slide.shapes.add_textbox(
        Emu(348016), Emu(6455949), Emu(11131033), Emu(215444)
    )
    tf = footer_txBox.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "PFs: Provident Funds; pps: Percentage Points"
    run.font.size = FOOTER_FONT_SIZE
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    return slide


def build_top10_aif_slide(prs, df):
    """Slide: Top 10 AIFs (7 columns)."""
    slide = add_blank_slide(prs)
    _add_title_textbox(slide, "Top 10 AIFs")
    _add_note_textbox(slide, "Shares in Lakhs")

    num_data_rows = min(len(df), 10)
    total_rows = 2 + num_data_rows
    num_cols = 7

    table_shape = slide.shapes.add_table(
        total_rows, num_cols,
        Emu(348016), Emu(1032476), Emu(11521440), Emu(2395728)
    )
    table = table_shape.table
    remove_table_styling(table)

    col_widths = [548640, 4663440, 1280160, 1280160, 1280160, 1280160, 1280160]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    _build_7col_header(table, CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE,
                       "Change in Holding\n(pps)")

    for r in range(2):
        for c in range(num_cols):
            set_cell_border(table.cell(r, c), border_color="FFFFFF")

    _fill_data_rows_7col(table, df, start_row=2, change_col_idx=6)

    # Footer
    footer_txBox = slide.shapes.add_textbox(
        Emu(348016), Emu(6455949), Emu(11131033), Emu(215444)
    )
    tf = footer_txBox.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "AIFs: Alternate Investment Funds; pps: Percentage Points"
    run.font.size = FOOTER_FONT_SIZE
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    return slide


def build_thankyou_slide(prs):
    """Build the 'Thank You' closing slide."""
    slide = add_blank_slide(prs)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x2B, 0x5C)  # Dark navy

    txBox = slide.shapes.add_textbox(
        Emu(4870309), Emu(2094717), Emu(2451383), Emu(569479)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Thank You"
    run.font.size = Pt(28)
    run.font.bold = False
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return slide


# â”€â”€ Main Orchestrator â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def generate_report(template_path: str, db_path: str, date_input: Optional[str], bu_id: int) -> tuple:
    """Library entry point to generate the PPT report.
    Returns: (pptx_bytes: bytes, display_date: str, report_data: dict)
    """
    global DB_PATH, TEMPLATE_PPT
    DB_PATH = db_path
    TEMPLATE_PPT = template_path

    conn = get_db_connection()
    bu_name = "Adani Portfolio"
    try:
        try:
            cur = conn.cursor()
            cur.execute("SELECT bu_name FROM bu_details WHERE bu_id = ?", (bu_id,))
            row = cur.fetchone()
            if row:
                bu_name = row[0]
        except Exception:
            pass

        date_range_label = date_input
        if not date_range_label or date_range_label == "latest":
            date_range_label = _pick_latest_daterange(conn, TABLE_NAMES["top_20_institutional"], bu_id=bu_id)

        # Keep date strings consistent across the deck.
        # DateRange label is stored as: "<previous> vs <current>".
        d1, d2 = "N/A", "N/A"  # d1=current, d2=previous
        if date_range_label and " vs " in date_range_label:
            prev_s, cur_s = date_range_label.split(" vs ")
            d2, d1 = prev_s.strip(), cur_s.strip()

        # Ensure global dates used by slide builders match the selected DateRange.
        # (These globals are referenced widely in existing helper functions.)
        global CURRENT_WEEK_DATE, PREVIOUS_WEEK_DATE
        if d1 != "N/A":
            CURRENT_WEEK_DATE = d1
        if d2 != "N/A":
            PREVIOUS_WEEK_DATE = d2

        df_institutional = fetch_table_data_latest(conn, TABLE_NAMES["top_20_institutional"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_buyers = fetch_table_data_latest(conn, TABLE_NAMES["top_20_buyers"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_sellers = fetch_table_data_latest(conn, TABLE_NAMES["top_20_sellers"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_fii_fpi = fetch_table_data_latest(conn, TABLE_NAMES["top_10_fii_fpi"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_mf_active = fetch_table_data_latest(conn, TABLE_NAMES["top_10_mf_active"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_mf_passive = fetch_table_data_latest(conn, TABLE_NAMES["top_10_mf_passive"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_insurance_pf = fetch_table_data_latest(conn, TABLE_NAMES["top_10_insurance_pf"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_aif = fetch_table_data_latest(conn, TABLE_NAMES["top_10_aif"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_entry = fetch_table_data_latest(conn, TABLE_NAMES["entry"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)
        df_exit = fetch_table_data_latest(conn, TABLE_NAMES["exit"], SCHEMA_NAME, date_range=date_range_label, bu_id=bu_id)

        import io
        prs = Presentation(template_path)

        # Restore missing template orchestration:
        # - Update Title slide date (and keep its layout)
        # - Update Table of Contents slide content
        try:
            build_title_slide(prs, bu_name=bu_name)
        except Exception:
            pass
        try:
            toc_items = [
                "Top 20 Institutional Shareholders",
                "Top 20 Buyers",
                "Top 20 Sellers",
                "New Entry / Exits",
                "Top 10 FIIs & FPIs",
                "Top 10 MFs",
                "Top 10 Insurance & PFs",
                "Top 10 AIFs",
            ]
            update_toc_on_template(prs, toc_items)
        except Exception:
            pass

        ok1 = update_institutional_table_on_template(prs, df_institutional)
        ok2 = update_buyers_table_on_template(prs, df_buyers)
        ok3 = update_sellers_table_on_template(prs, df_sellers)
        ok4 = update_new_entry_exits_on_template(prs, df_entry, df_exit)
        ok5 = update_fii_fpi_table_on_template(prs, df_fii_fpi)
        ok6 = update_mf_tables_on_template(prs, df_mf_active, df_mf_passive)
        ok7 = update_insurance_pf_table_on_template(prs, df_insurance_pf)
        ok8 = update_aif_table_on_template(prs, df_aif)
        print(f"[PPT][GEN] Updates: Inst={ok1}, Buy={ok2}, Sell={ok3}, EntryExit={ok4}, FII={ok5}, MF={ok6}, Ins={ok7}, AIF={ok8}")

        raw_delete = (_get_env("WSHP_DELETE_SLIDES", "3,12,13") or "").strip()
        if raw_delete:
            try:
                def _norm_txt(s: str) -> str:
                    return " ".join(str(s or "").split()).strip().lower()

                def _slide_title_text(slide) -> str:
                    try:
                        for shp in getattr(slide, "shapes", []):
                            try:
                                if getattr(shp, "has_text_frame", False):
                                    t = _norm_txt(shp.text_frame.text)
                                    if t:
                                        return t
                            except Exception:
                                continue
                    except Exception:
                        pass
                    return ""

                protected_keywords = [
                    "new entry",
                    "new entries",
                    "exit",
                    "exits",
                    "entry / exits",
                    "entry/exits",
                    "entry",
                ]

                one_based = [int(x.strip()) for x in raw_delete.split(",") if x.strip()]
                cand = [i - 1 for i in one_based if i > 0]
                # Never delete slide #7 (1-based), i.e. index 6 (0-based)
                cand = [i for i in cand if i != 6]
                safe = []
                for idx in cand:
                    try:
                        if idx < 0 or idx >= len(prs.slides):
                            continue
                        title_txt = _slide_title_text(prs.slides[idx])
                        if any(k in title_txt for k in protected_keywords):
                            continue
                        safe.append(idx)
                    except Exception:
                        continue
                if safe:
                    _delete_slides(prs, safe)
            except Exception:
                pass

        for slide in prs.slides:
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        txt = shape.text_frame.text or ""
                        if "Portfolio" in txt or "PORTFOLIO" in txt or "Adani" in txt:
                            new_txt = txt.replace("Adani Portfolio", bu_name).replace("ADANI PORTFOLIO", bu_name.upper())
                            if txt.strip() in ["Adani Portfolio", "ADANI PORTFOLIO"]:
                                new_txt = bu_name
                            if new_txt != txt:
                                shape.text_frame.text = new_txt
                        if "[Date]" in txt or "(Date)" in txt or "DATE" in txt:
                            shape.text_frame.text = txt.replace("[Date]", d1).replace("(Date)", d1).replace("DATE", d1)
            except Exception:
                pass

        report_data = {
            "date_range_label": date_range_label,
            "bu_id": bu_id,
            "bu_name": bu_name,
        }

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue(), d1, report_data
    finally:
        conn.close()


def main():
    print(f"\n{'=' * 65}")
    print("  POWERPOINT REPORT GENERATOR (SQLite)")
    print(f"{'=' * 65}")

    bu_id = _get_bu_id()
    date_range = _get_date_range_label()
    pptx_bytes, display_date, _ = generate_report(TEMPLATE_PPT, DB_PATH, date_range, bu_id)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), OUTPUT_PPT_FILENAME)
    with open(output_path, "wb") as f:
        f.write(pptx_bytes)

    print("  OK SUCCESS!")
    print(f"   Output: {output_path}")
    print(f"{'=' * 65}\n")


if __name__ == "__main__":
    main()
